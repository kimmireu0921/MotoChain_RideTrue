#!/usr/bin/env python3
"""
MotoChain & RideTrue — whitepaper-based presentation builder.
15 slides with matplotlib charts, flow diagrams, and clean dark theme.
"""

import io
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────
BG_H  = '#060c16'
BG2_H = '#0d1624'
BLU_H = '#4aa3ff'
GRN_H = '#4ee59a'
AMB_H = '#ffbf56'
CYN_H = '#27d8d8'
RED_H = '#ff6b6b'
TXT_H = '#e8eaf0'
MUT_H = '#8892a0'
LIN_H = '#1e2a3a'

def _rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

BG_R = _rgb(BG_H);  TXT_R = _rgb(TXT_H); MUT_R = _rgb(MUT_H)
BLU_R = _rgb(BLU_H); GRN_R = _rgb(GRN_H); AMB_R = _rgb(AMB_H)
CYN_R = _rgb(CYN_H); RED_R = _rgb(RED_H); LIN_R = _rgb(LIN_H)

# Slide dimensions
SW = Inches(13.333)
SH = Inches(7.5)
ML = Inches(0.48)   # left margin
DPI = 150

# Matplotlib defaults
plt.rcParams.update({
    'figure.facecolor': BG_H, 'axes.facecolor': BG2_H,
    'savefig.facecolor': BG_H, 'text.color': TXT_H,
    'axes.labelcolor': TXT_H, 'xtick.color': MUT_H, 'ytick.color': MUT_H,
    'axes.edgecolor': LIN_H, 'grid.color': LIN_H, 'font.family': 'sans-serif',
})

# ── matplotlib helpers ────────────────────────────────────────────────────

def fig2stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=DPI, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    plt.close(fig)
    buf.seek(0)
    return buf

def hflow(steps, fig_w=12.4, fig_h=2.6):
    """Generic horizontal flow-chart."""
    fig, ax = plt.subplots(figsize=(fig_w, fig_h), facecolor=BG_H)
    ax.set_facecolor(BG_H); ax.axis('off')
    n = len(steps)
    bw = 0.16; gap = 0.04
    total = n*bw + (n-1)*gap
    sx = (1.0 - total) / 2
    for i, (num, label, body, col) in enumerate(steps):
        x = sx + i*(bw+gap)
        patch = FancyBboxPatch((x, 0.1), bw, 0.78,
                               boxstyle="round,pad=0.02",
                               fc=col+'22', ec=col, lw=2,
                               transform=ax.transAxes, clip_on=False)
        ax.add_patch(patch)
        ax.text(x+bw/2, 0.77, num,   transform=ax.transAxes,
                ha='center', va='center', fontsize=8.5, fontweight='bold', color=col)
        ax.text(x+bw/2, 0.58, label, transform=ax.transAxes,
                ha='center', va='center', fontsize=9,   fontweight='bold', color=TXT_H)
        ax.text(x+bw/2, 0.30, body,  transform=ax.transAxes,
                ha='center', va='center', fontsize=7.2, color=MUT_H, linespacing=1.4)
        if i < n-1:
            ax.annotate('', xy=(x+bw+gap-0.006, 0.5), xytext=(x+bw+0.005, 0.5),
                        xycoords='axes fraction', textcoords='axes fraction',
                        arrowprops=dict(arrowstyle='->', color=MUT_H, lw=1.8))
    return fig2stream(fig)

# ── CHARTS ───────────────────────────────────────────────────────────────

def make_vin_flow():
    steps = [
        ('01', 'MINT',     'ERC-721 NFT\ncreated', BLU_H),
        ('02', 'REGISTER', 'Metadata\nattached', CYN_H),
        ('03', 'LOG',      'Mechanic + owner\nco-sign records', GRN_H),
        ('04', 'TRANSFER', 'Ownership audited\non-chain', AMB_H),
        ('05', 'QUERY',    'Insurers & buyers\naccess history', RED_H),
    ]
    return hflow(steps)

def make_zk_flow():
    steps = [
        ('01', 'COLLECT', 'GPS · accel\ngyro · ECU', BLU_H),
        ('02', 'COMPUTE', 'Off-chain\nbehavior score', CYN_H),
        ('03', 'PROVE',   'ZK proof\ngenerated', AMB_H),
        ('04', 'SUBMIT',  'Proof verified\non-chain', GRN_H),
        ('05', 'UPDATE',  'Score & premium\nupdated', GRN_H),
    ]
    return hflow(steps)

def make_claim_flow():
    n = 6
    bw = 0.135; gap = 0.028
    total = n*bw+(n-1)*gap; sx = (1-total)/2
    steps = [
        ('INCIDENT',  'Accident,\ntheft or damage',      RED_H),
        ('SUBMIT',    'Rider submits\nvia app',           BLU_H),
        ('ORACLE\nVERIFY','Multi-source\nevidence',       CYN_H),
        ('CONSENSUS', 'Validators reach\nquorum',         AMB_H),
        ('PAYOUT',    'USDC released\nfrom pool',         GRN_H),
        ('REBALANCE', 'Reserves &\npremiums reset',       GRN_H),
    ]
    fig, ax = plt.subplots(figsize=(12.4, 2.4), facecolor=BG_H)
    ax.set_facecolor(BG_H); ax.axis('off')
    for i, (label, body, col) in enumerate(steps):
        x = sx + i*(bw+gap)
        patch = FancyBboxPatch((x, 0.08), bw, 0.84,
                               boxstyle="round,pad=0.02",
                               fc=col+'22', ec=col, lw=1.8,
                               transform=ax.transAxes, clip_on=False)
        ax.add_patch(patch)
        ax.text(x+bw/2, 0.72, label, transform=ax.transAxes,
                ha='center', va='center', fontsize=8.5, fontweight='bold',
                color=col, linespacing=1.3)
        ax.text(x+bw/2, 0.28, body,  transform=ax.transAxes,
                ha='center', va='center', fontsize=7, color=TXT_H, linespacing=1.4)
        if i < n-1:
            ax.annotate('', xy=(x+bw+gap-0.005, 0.5), xytext=(x+bw+0.003, 0.5),
                        xycoords='axes fraction', textcoords='axes fraction',
                        arrowprops=dict(arrowstyle='->', color=MUT_H, lw=1.6))
    return fig2stream(fig)

def make_revenue_pie():
    fig, ax = plt.subplots(figsize=(5.6, 4.8),
                           subplot_kw=dict(aspect='equal'), facecolor=BG_H)
    ax.set_facecolor(BG_H)
    labels = ['Insurance Pool\nFee', 'API / Data\nAccess', 'VIN\nRegistration',
              'Ownership\nTransfer', 'Service\nLogging', 'Claims\nProcessing']
    sizes  = [44, 22, 22, 7, 4, 1]
    colors = [BLU_H, GRN_H, CYN_H, AMB_H, RED_H, MUT_H]
    wedges, _ = ax.pie(sizes, colors=colors, startangle=90,
                       wedgeprops=dict(linewidth=2, edgecolor=BG_H),
                       explode=[0.03]*len(sizes))
    legend_labels = [f'{l.replace(chr(10)," ")}  {s}%' for l,s in zip(labels,sizes)]
    ax.legend(wedges, legend_labels, loc='lower center', bbox_to_anchor=(0.5,-0.18),
              ncol=2, fontsize=8, frameon=False, labelcolor=TXT_H)
    ax.set_title('Year 1 Revenue Mix', color=TXT_H, fontsize=11, fontweight='bold', pad=8)
    fig.tight_layout()
    return fig2stream(fig)

def make_token_donut():
    fig, ax = plt.subplots(figsize=(5.6, 4.8),
                           subplot_kw=dict(aspect='equal'), facecolor=BG_H)
    ax.set_facecolor(BG_H)
    labels = ['Community &\nEcosystem 30%', 'Team &\nFounders 18%',
              'Treasury / DAO 15%', 'Pool Subsidy 15%',
              'Investors\nPrivate 9%', 'Investors\nSeed 8%', 'Liquidity 5%']
    sizes  = [30, 18, 15, 15, 9, 8, 5]
    colors = [BLU_H, GRN_H, CYN_H, AMB_H, RED_H, MUT_H, '#a855f7']
    wedges, _ = ax.pie(sizes, colors=colors, startangle=90,
                       wedgeprops=dict(width=0.55, linewidth=2, edgecolor=BG_H))
    ax.text(0, 0, '1 B\n$MOTO', ha='center', va='center',
            fontsize=13, fontweight='bold', color=TXT_H, linespacing=1.6)
    ax.legend(wedges, labels, loc='lower center', bbox_to_anchor=(0.5,-0.18),
              ncol=2, fontsize=7.5, frameon=False, labelcolor=TXT_H)
    ax.set_title('$MOTO Token Allocation', color=TXT_H, fontsize=11, fontweight='bold', pad=8)
    fig.tight_layout()
    return fig2stream(fig)

def make_premium_pie():
    fig, ax = plt.subplots(figsize=(5.6, 4.8),
                           subplot_kw=dict(aspect='equal'), facecolor=BG_H)
    ax.set_facecolor(BG_H)
    labels = ['Liquid Reserve\n25%', 'Catastrophe\nReserve 20%',
              'Safe-Rider\nRebates 18%', 'Expected\nClaims 12%',
              'Treasury\n10%', 'Reinsurance\n10%', 'Protocol\nFee 5%']
    sizes  = [25, 20, 18, 12, 10, 10, 5]
    colors = [GRN_H, BLU_H, CYN_H, RED_H, AMB_H, MUT_H, '#a855f7']
    wedges, _ = ax.pie(sizes, colors=colors, startangle=90,
                       wedgeprops=dict(linewidth=2, edgecolor=BG_H),
                       explode=[0.03]*len(sizes))
    ax.legend(wedges, labels, loc='lower center', bbox_to_anchor=(0.5,-0.18),
              ncol=2, fontsize=8, frameon=False, labelcolor=TXT_H)
    ax.set_title('Where Every Premium Goes', color=TXT_H, fontsize=11, fontweight='bold', pad=8)
    fig.tight_layout()
    return fig2stream(fig)

def make_pnl_bars():
    fig, ax = plt.subplots(figsize=(9.5, 4.2), facecolor=BG_H)
    ax.set_facecolor(BG2_H)
    years   = ['Year 1', 'Year 2', 'Year 3']
    revenue = [275,   1515,  6448]
    opex    = [-1800, -2500, -3400]
    net     = [-1525, -985,   3048]
    x = np.arange(3); w = 0.24
    ax.bar(x-w,   revenue, w, label='Revenue (₩M)', color=GRN_H, alpha=0.9, zorder=3)
    ax.bar(x,     opex,    w, label='OpEx (₩M)',    color=RED_H, alpha=0.9, zorder=3)
    ax.bar(x+w,   net,     w, label='Net (₩M)',     color=BLU_H, alpha=0.9, zorder=3)
    for i,(r,e,n) in enumerate(zip(revenue,opex,net)):
        ax.text(i-w, r+80, f'₩{r:,}', ha='center', fontsize=7.5, color=GRN_H, fontweight='bold')
        ax.text(i,   e-80, f'₩{e:,}', ha='center', fontsize=7.5, color=RED_H, fontweight='bold', va='top')
        va = 'bottom' if n>=0 else 'top'
        ax.text(i+w, n+(80 if n>=0 else -80), f'₩{n:,}', ha='center',
                fontsize=7.5, color=BLU_H, fontweight='bold', va=va)
    ax.axhline(0, color=MUT_H, lw=0.8, ls='--', zorder=2)
    ax.annotate('Break-even\nYear 3 →', xy=(2+w, 3048), xytext=(2.15, 2200),
                arrowprops=dict(arrowstyle='->', color=GRN_H, lw=1.5),
                color=GRN_H, fontsize=9, fontweight='bold')
    ax.set_xticks(x); ax.set_xticklabels(years, fontsize=11, color=TXT_H)
    ax.set_ylabel('₩ Millions', color=MUT_H, fontsize=9)
    ax.tick_params(axis='y', labelcolor=MUT_H, labelsize=8)
    ax.spines[['top','right']].set_visible(False)
    ax.spines[['left','bottom']].set_color(LIN_H)
    ax.grid(axis='y', alpha=0.2, color=MUT_H, ls='--', zorder=1)
    ax.legend(fontsize=9, frameon=False, labelcolor=TXT_H, loc='upper left')
    ax.set_title('3-Year Stressed P&L (15% claim frequency · ₩2.5M avg severity)',
                 color=TXT_H, fontsize=10, pad=8)
    fig.tight_layout()
    return fig2stream(fig)

def make_lp_bars():
    fig, ax = plt.subplots(figsize=(8.5, 3.8), facecolor=BG_H)
    ax.set_facecolor(BG2_H)
    years = ['Year 1', 'Year 2', 'Year 3']
    foundation    = [1500, 1500,  1500]
    defi          = [750,  4000, 15000]
    institutional = [500,  3000, 12000]
    strategic     = [250,  2000,  8000]
    x = np.arange(3); w = 0.45
    b1 = ax.bar(x, foundation, w, label='Foundation', color=BLU_H, alpha=0.9)
    b2 = ax.bar(x, defi, w, bottom=foundation, label='DeFi LPs', color=GRN_H, alpha=0.9)
    bot3 = [foundation[i]+defi[i] for i in range(3)]
    b3 = ax.bar(x, institutional, w, bottom=bot3, label='Institutional', color=AMB_H, alpha=0.9)
    bot4 = [bot3[i]+institutional[i] for i in range(3)]
    b4 = ax.bar(x, strategic, w, bottom=bot4, label='Strategic', color=CYN_H, alpha=0.9)
    totals = [3000, 10500, 36500]
    for i,t in enumerate(totals):
        ax.text(i, bot4[i]+strategic[i]+300, f'Total\n₩{t:,}M',
                ha='center', fontsize=8.5, color=TXT_H, fontweight='bold', linespacing=1.4)
    ax.set_xticks(x); ax.set_xticklabels(years, fontsize=11, color=TXT_H)
    ax.set_ylabel('₩ Millions', color=MUT_H, fontsize=9)
    ax.tick_params(axis='y', labelcolor=MUT_H, labelsize=8)
    ax.spines[['top','right']].set_visible(False)
    ax.spines[['left','bottom']].set_color(LIN_H)
    ax.grid(axis='y', alpha=0.2, color=MUT_H, ls='--')
    ax.legend(fontsize=8.5, frameon=False, labelcolor=TXT_H, loc='upper left')
    ax.set_title('LP Capital Plan (₩ Millions)', color=TXT_H, fontsize=10.5, pad=8)
    fig.tight_layout()
    return fig2stream(fig)

def make_mechanic_defense():
    fig, ax = plt.subplots(figsize=(12.4, 3.6), facecolor=BG_H)
    ax.set_facecolor(BG_H); ax.axis('off')
    layers = [
        ('PREVENTION', GRN_H, [
            ('Multi-party co-signing', 'Mechanic AND owner must both sign every record.\nNo one can fabricate a repair unilaterally.'),
            ('Hardware-attested telemetry', 'OBD-II reads mileage directly from the bike.\nNo room to fudge the odometer.'),
        ]),
        ('DETECTION', AMB_H, [
            ('Bounty market', 'Anyone who spots a fake record and proves it\nearns a share of the slashed stake.'),
            ('Re-verify at monetization', 'History is re-checked at every resale, loan,\nor insurance event — when fraud would pay off.'),
        ]),
        ('ENFORCEMENT', RED_H, [
            ('Tiered staking', 'More write volume = more stake at risk.\nLarge-scale fraud cannot cap its downside.'),
            ('Escrow delay 30–90 days', 'Fees sit in escrow. Any dispute claws fees back\nand slashes stake proportional to entry volume.'),
        ]),
    ]
    cw = 0.295; gap = 0.035
    total = 3*cw + 2*gap; sx = (1-total)/2
    for i, (title, col, items) in enumerate(layers):
        x = sx + i*(cw+gap)
        # header
        hdr = FancyBboxPatch((x, 0.82), cw, 0.15,
                             boxstyle="round,pad=0.01",
                             fc=col+'33', ec=col, lw=2,
                             transform=ax.transAxes, clip_on=False)
        ax.add_patch(hdr)
        ax.text(x+cw/2, 0.895, title, transform=ax.transAxes,
                ha='center', va='center', fontsize=10.5, fontweight='bold', color=col)
        # body
        body_box = FancyBboxPatch((x, 0.02), cw, 0.78,
                                  boxstyle="round,pad=0.01",
                                  fc=col+'0d', ec=col+'55', lw=1,
                                  transform=ax.transAxes, clip_on=False)
        ax.add_patch(body_box)
        y = 0.74
        for j, (item_title, item_body) in enumerate(items):
            if j > 0: y -= 0.08
            ax.text(x+0.015, y, item_title, transform=ax.transAxes,
                    ha='left', va='top', fontsize=9.5, fontweight='bold', color=col)
            y -= 0.125
            ax.text(x+0.015, y, item_body, transform=ax.transAxes,
                    ha='left', va='top', fontsize=8, color=TXT_H, linespacing=1.45)
            y -= 0.15
    ax.text(0.5, 0.0, 'Each layer covers what the others miss — none works alone — all three make fraud unprofitable',
            transform=ax.transAxes, ha='center', va='bottom', fontsize=8, color=MUT_H, style='italic')
    return fig2stream(fig)

def make_competitive():
    fig, ax = plt.subplots(figsize=(12.0, 3.8), facecolor=BG_H)
    ax.set_facecolor(BG2_H)
    comps    = ['Centralized History\n(CARFAX, Korea Auto)', 'Traditional Insurers\n(Samsung, Hyundai M&F)',
                'DeFi Insurance\n(Nexus Mutual, InsurAce)', 'MotoChain & RideTrue ←']
    criteria = ['Portable\nIdentity', 'Behavior-\nBased Price', 'On-Chain\nClaims',
                'Economic\nEnforcement', 'Consumer\nUX', 'Real-World\nData Moat']
    scores = [
        [0, 0, 0, 0, 1, 2],
        [0, 1, 0, 0, 2, 1],
        [1, 0, 2, 1, 0, 0],
        [2, 2, 2, 2, 2, 2],
    ]
    sym   = {0:'✗', 1:'∼', 2:'✓'}
    fc    = {0: RED_H+'44', 1: AMB_H+'44', 2: GRN_H+'44'}
    tc    = {0: RED_H,      1: AMB_H,      2: GRN_H}
    for ri, row in enumerate(scores):
        for ci, s in enumerate(row):
            rect = mpatches.FancyBboxPatch((ci+0.06, ri+0.06), 0.88, 0.88,
                                           boxstyle="round,pad=0.05",
                                           fc=fc[s], ec='none')
            ax.add_patch(rect)
            ax.text(ci+0.5, ri+0.5, sym[s], ha='center', va='center',
                    fontsize=17, fontweight='bold', color=tc[s])
    # highlight MotoChain row
    hl = mpatches.FancyBboxPatch((-0.06, 3.02), len(criteria)+0.12, 0.96,
                                  boxstyle="round,pad=0.02",
                                  fc='none', ec=BLU_H, lw=2.5)
    ax.add_patch(hl)
    ax.set_xlim(0, len(criteria)); ax.set_ylim(0, len(comps))
    ax.set_xticks([i+0.5 for i in range(len(criteria))])
    ax.set_xticklabels(criteria, fontsize=8.5, color=TXT_H, linespacing=1.35)
    ax.set_yticks([i+0.5 for i in range(len(comps))])
    ax.set_yticklabels(comps, fontsize=9, color=TXT_H)
    ax.tick_params(length=0)
    ax.spines[:].set_visible(False)
    ax.set_title('Feature Comparison Matrix', color=TXT_H, fontsize=10.5, pad=8)
    # legend
    for label, col in [('✓ Yes', GRN_H), ('∼ Partial', AMB_H), ('✗ No', RED_H)]:
        ax.text(0, -0.3, '', color=col)  # placeholder
    from matplotlib.lines import Line2D
    legend_els = [Line2D([0],[0], marker='s', color='w', markerfacecolor=GRN_H, markersize=10, label='✓ Yes'),
                  Line2D([0],[0], marker='s', color='w', markerfacecolor=AMB_H, markersize=10, label='∼ Partial'),
                  Line2D([0],[0], marker='s', color='w', markerfacecolor=RED_H, markersize=10, label='✗ No')]
    ax.legend(handles=legend_els, loc='upper right', fontsize=8.5, frameon=False,
              labelcolor=TXT_H, bbox_to_anchor=(1, -0.08), ncol=3)
    fig.tight_layout()
    return fig2stream(fig)

def make_security_timeline():
    fig, ax = plt.subplots(figsize=(12.4, 3.4), facecolor=BG_H)
    ax.set_facecolor(BG_H); ax.axis('off')
    phases = [
        ('Pre-Launch\nEngineering Gate', BLU_H, [
            '≥90% unit test coverage', 'Invariant & fuzz testing',
            'Zero unresolved critical findings', 'Internal threat model review']),
        ('External\nAudits', CYN_H, [
            '2 independent smart-contract audits',
            'OpenZeppelin / Trail of Bits / Spearbit',
            'Formal verification of fund flows',
            'All reports published in full']),
        ('Bug Bounty\n(Immunefi)', AMB_H, [
            'Critical: up to $250,000',
            'High: up to $75,000',
            'Medium: up to $15,000',
            'Live for lifetime of protocol']),
        ('Operational\nSecurity', GRN_H, [
            '5-of-9 multi-sig treasury',
            '48-hour time-locked upgrades',
            'Emergency pause (narrow scope)',
            'Annual third-party pen tests']),
    ]
    n = len(phases); cw = 0.215; gap = 0.03
    total = n*cw + (n-1)*gap; sx = (1-total)/2
    ax.axhline(0.74, xmin=sx, xmax=sx+total, color=MUT_H, lw=1.5)
    for i, (title, col, pts) in enumerate(phases):
        x = sx + i*(cw+gap); cx = x + cw/2
        ax.plot(cx, 0.74, 'o', color=col, ms=13, transform=ax.transAxes, zorder=4)
        ax.text(cx, 0.85, title, transform=ax.transAxes,
                ha='center', va='bottom', fontsize=8.5, fontweight='bold',
                color=col, linespacing=1.4)
        body = FancyBboxPatch((x, 0.02), cw, 0.60,
                              boxstyle="round,pad=0.01",
                              fc=col+'18', ec=col+'55', lw=1,
                              transform=ax.transAxes, clip_on=False)
        ax.add_patch(body)
        for j, pt in enumerate(pts):
            ax.text(x+0.012, 0.57-j*0.125, f'• {pt}',
                    transform=ax.transAxes, ha='left', va='top', fontsize=7.5, color=TXT_H)
    return fig2stream(fig)

def make_ecosystem():
    fig, ax = plt.subplots(figsize=(7.5, 5.5), facecolor=BG_H)
    ax.set_facecolor(BG_H)
    ax.set_xlim(-1.3, 1.3); ax.set_ylim(-1.15, 1.15); ax.axis('off')
    # hub
    hub = plt.Circle((0,0), 0.27, color=BLU_H+'33', ec=BLU_H, lw=3, zorder=3)
    ax.add_patch(hub)
    ax.text(0, 0.05, 'MotoChain', ha='center', va='center',
            fontsize=9.5, fontweight='bold', color=BLU_H, zorder=4)
    ax.text(0, -0.10, '& RideTrue', ha='center', va='center',
            fontsize=8.5, color=BLU_H, zorder=4)
    spokes = [
        ('Riders',    0,     0.88,  GRN_H, 'Premium + safe-rider\nrebates'),
        ('Mechanics', -0.84, 0.48,  AMB_H, 'Stake + write access\n+ fee share'),
        ('Dealers',    0.84, 0.48,  AMB_H, 'Verified history\nat point of sale'),
        ('Insurers',  -0.84,-0.48,  CYN_H, 'Behavior data\nfor underwriting'),
        ('Lenders',    0,   -0.88,  CYN_H, 'Verified VIN history\nfor loan decisions'),
        ('API Users',  0.84,-0.48,  CYN_H, 'B2B subscription\ndata access'),
    ]
    for label, sx, sy, col, note in spokes:
        norm = (sx**2+sy**2)**0.5
        hx = sx/norm*0.27; hy = sy/norm*0.27
        ax.annotate('', xy=(sx*0.60, sy*0.60), xytext=(hx, hy),
                    arrowprops=dict(arrowstyle='->', color=col+'bb', lw=1.6))
        c = plt.Circle((sx*0.76, sy*0.76), 0.155, color=col+'22', ec=col, lw=2, zorder=3)
        ax.add_patch(c)
        ax.text(sx*0.76, sy*0.76+0.02, label, ha='center', va='center',
                fontsize=9, fontweight='bold', color=col, zorder=4)
        # note slightly outside
        nx = sx*1.05; ny = sy*1.05
        ax.text(nx, ny, note, ha='center', va='center',
                fontsize=6.8, color=MUT_H, linespacing=1.35)
    return fig2stream(fig)

# ── PPTX helpers ─────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def add_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_R
    return s

def tb(slide, text, l, t, w, h, size=11, bold=False,
       color=TXT_R, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run()
    r.text = text; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return box

def add_para(tf, text, size=10, bold=False, color=TXT_R, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

def box_shape(slide, l, t, w, h, fill_hex, line_hex=None, line_w=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = _rgb(fill_hex)
    if line_hex:
        s.line.color.rgb = _rgb(line_hex)
        if line_w: s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s

def add_img(slide, stream, l, t, w, h=None):
    stream.seek(0)
    if h:
        slide.shapes.add_picture(stream, l, t, w, h)
    else:
        slide.shapes.add_picture(stream, l, t, w)

def eyebrow(slide, text):
    tb(slide, text.upper(), ML, Inches(0.27), SW-Inches(0.96), Inches(0.3),
       size=8, bold=True, color=BLU_R)

def heading(slide, text, size=28, top=Inches(0.55), color=TXT_R):
    tb(slide, text, ML, top, SW-Inches(0.96), Inches(0.62),
       size=size, bold=True, color=color)

def pagenum(slide, n, total):
    tb(slide, f'{n} / {total}',
       SW-Inches(1.05), SH-Inches(0.38), Inches(0.95), Inches(0.28),
       size=8, color=MUT_R, align=PP_ALIGN.RIGHT)

# ── SLIDE BUILDERS ────────────────────────────────────────────────────────

def s01_title(slide):
    # Big title
    tb(slide, 'MotoChain & RideTrue',
       Inches(0.7), Inches(1.6), SW-Inches(4.4), Inches(1.3),
       size=42, bold=True, color=TXT_R)
    tb(slide, 'A Verifiable Identity & Behavior-Based\nInsurance Protocol for Motorcycles',
       Inches(0.7), Inches(3.0), SW-Inches(4.4), Inches(1.0),
       size=17, color=MUT_R)
    # accent rule
    box_shape(slide, Inches(0.7), Inches(4.1), Inches(6.5), Pt(3), BLU_H)
    tb(slide, 'Mireu Kim  ·  Enxhi Topalli',
       Inches(0.7), Inches(4.3), Inches(7), Inches(0.4),
       size=13, color=BLU_R)
    tb(slide, 'CS495 DeFi  ·  May 2025',
       Inches(0.7), Inches(4.75), Inches(7), Inches(0.35),
       size=11, color=MUT_R)
    # right-side KPI boxes
    kpis = [('$120 B+','Global Motorcycle Market', BLU_H),
            ('$70 B+', 'Motorcycle Insurance Market', GRN_H),
            ('Base L2', 'Launch Network', CYN_H)]
    for i,(kpi,label,col) in enumerate(kpis):
        x = Inches(9.8); y = Inches(2.1 + i*1.55)
        box_shape(slide, x, y, Inches(3.1), Inches(1.3), LIN_H, col, 1.5)
        tb(slide, kpi,   x+Inches(0.18), y+Inches(0.1),  Inches(2.7), Inches(0.6),
           size=24, bold=True, color=_rgb(col))
        tb(slide, label, x+Inches(0.18), y+Inches(0.72), Inches(2.7), Inches(0.42),
           size=9.5, color=MUT_R)

def s02_problem(slide):
    eyebrow(slide, 'The Problem')
    heading(slide, 'Two Broken Markets — One Root Cause')
    t = Inches(1.2)
    panels = [
        ('$120 B+', 'Global Motorcycle Market', BLU_H,
         'Buyers cannot verify accident history, mileage integrity, or service quality. '
         'Honest sellers are discounted because the market assumes risk — the classic "market for lemons."'),
        ('$70 B+', 'Insurance Market', AMB_H,
         'Insurers rely on age and demographics instead of real riding behavior. '
         'Safe riders subsidize risky ones. Premiums do not reflect actual risk.'),
    ]
    bw = (SW - ML - Inches(0.45) - Inches(0.25)) / 2
    for i,(kpi,label,col,desc) in enumerate(panels):
        x = ML + i*(bw+Inches(0.25))
        box_shape(slide, x, t, bw, Inches(4.6), LIN_H, col, 1.5)
        tb(slide, kpi,   x+Inches(0.2), t+Inches(0.15), bw-Inches(0.25), Inches(0.85),
           size=36, bold=True, color=_rgb(col))
        tb(slide, label, x+Inches(0.2), t+Inches(1.0),  bw-Inches(0.25), Inches(0.45),
           size=14, bold=True, color=TXT_R)
        tb(slide, desc,  x+Inches(0.2), t+Inches(1.55), bw-Inches(0.35), Inches(2.6),
           size=11, color=MUT_R)
    # Root cause bar
    box_shape(slide, ML, t+Inches(4.85), SW-ML-Inches(0.45), Inches(0.8), '#0d1220', RED_H, 1.5)
    tb(slide, 'Root cause:  No shared, tamper-proof, portable record of real-world events',
       ML+Inches(0.25), t+Inches(5.0), SW-ML-Inches(0.9), Inches(0.5),
       size=13, bold=True, color=RED_R)

def s03_solution(slide):
    eyebrow(slide, 'The Solution')
    heading(slide, 'A Dual-Protocol DeFi System on Base (Ethereum L2)')
    t = Inches(1.25)
    cols = [
        ('MotoChain', 'Identity Layer', BLU_H,
         ['ERC-721 VIN-linked NFT for every motorcycle',
          'Mechanic + owner co-signed service records',
          'Tamper-proof ownership transfer history',
          'Cryptographic proof accessible to any buyer',
          'Portable identity — follows the bike, not the platform']),
        ('RideTrue', 'Insurance Layer', GRN_H,
         ['Behavior-based premiums using ZK proofs',
          'Pooled stablecoin (USDC) insurance fund',
          'Automated on-chain claim processing',
          'Safe-rider quarterly rebates for low-risk behavior',
          'Composable DeFi reserve yield (Aave / Compound)']),
    ]
    bw = (SW - ML - Inches(0.45) - Inches(0.25)) / 2
    for i,(title,sub,col,pts) in enumerate(cols):
        x = ML + i*(bw+Inches(0.25))
        box_shape(slide, x, t, bw, Inches(5.55), LIN_H, col, 1.5)
        tb(slide, title, x+Inches(0.2), t+Inches(0.15), bw, Inches(0.55),
           size=26, bold=True, color=_rgb(col))
        tb(slide, sub,   x+Inches(0.2), t+Inches(0.72), bw, Inches(0.3),
           size=11, italic=True, color=MUT_R)
        bxt = slide.shapes.add_textbox(x+Inches(0.2), t+Inches(1.1), bw-Inches(0.4), Inches(4.1))
        tf = bxt.text_frame; tf.word_wrap = True
        for j,pt in enumerate(pts):
            p = tf.paragraphs[0] if j==0 else tf.add_paragraph()
            p.space_before = Pt(7)
            r = p.add_run(); r.text = f'• {pt}'
            r.font.size = Pt(11); r.font.color.rgb = TXT_R
    tb(slide, '⛓  Built on Base L2 — lower fees · faster finality · Ethereum security',
       ML, t+Inches(5.72), SW-ML-Inches(0.45), Inches(0.35),
       size=10, italic=True, color=CYN_R)

def s04_vin(slide, img):
    eyebrow(slide, 'Technical Architecture — MotoChain')
    heading(slide, 'VIN Token Lifecycle')
    tb(slide, 'Each motorcycle is anchored by an ERC-721 NFT. Every event in its life is co-signed and permanently locked on-chain.',
       ML, Inches(1.22), SW-ML-Inches(0.45), Inches(0.35), size=10.5, italic=True, color=MUT_R)
    add_img(slide, img, ML, Inches(1.65), SW-ML-Inches(0.45), Inches(2.1))
    props = [
        ('Tamper-Resistant', 'Once written, no single actor can alter the record', BLU_H),
        ('Multi-Party Signed', 'Mechanic + owner both sign every service entry', GRN_H),
        ('Portable', 'History follows the bike across platforms and borders', CYN_H),
        ('Composable', 'Open API for insurers, lenders, and marketplaces', AMB_H),
    ]
    pw = (SW-ML-Inches(0.45)) / 4
    for i,(title,desc,col) in enumerate(props):
        x = ML + i*pw
        box_shape(slide, x, Inches(3.95), pw-Inches(0.12), Inches(2.75), LIN_H, col, 1.2)
        tb(slide, title, x+Inches(0.14), Inches(4.1),  pw-Inches(0.28), Inches(0.38),
           size=11, bold=True, color=_rgb(col))
        tb(slide, desc,  x+Inches(0.14), Inches(4.52), pw-Inches(0.28), Inches(2.0),
           size=9.5, color=TXT_R)

def s05_zk(slide, img):
    eyebrow(slide, 'Technical Architecture — RideTrue')
    heading(slide, 'Scoring Behavior Without Exposing Your Location')
    tb(slide, 'Zero-knowledge proofs let the protocol verify safe riding without ever storing raw GPS or sensor data on-chain.',
       ML, Inches(1.22), SW-ML-Inches(0.45), Inches(0.35), size=10.5, italic=True, color=MUT_R)
    add_img(slide, img, ML, Inches(1.65), SW-ML-Inches(0.45), Inches(2.1))
    t = Inches(3.9)
    tb(slide, 'What ZK Proves',       ML,            t, Inches(5.8), Inches(0.35), size=12, bold=True, color=GRN_R)
    tb(slide, 'What ZK Cannot Prove', Inches(7.0),   t, Inches(5.9), Inches(0.35), size=12, bold=True, color=RED_R)
    proves = ['Rider met speed and braking thresholds',
              'Computation was applied to some input',
              'Proof is valid per the circuit definition']
    cannot = ['Where the rider actually traveled',
              'That input data reflects real riding',
              'That the recording device was authentic']
    for i,(p,c) in enumerate(zip(proves,cannot)):
        tb(slide, f'✓  {p}', ML,          t+Inches(0.42+i*0.52), Inches(6.2), Inches(0.45), size=10, color=TXT_R)
        tb(slide, f'✗  {c}', Inches(7.0), t+Inches(0.42+i*0.52), Inches(6.0), Inches(0.45), size=10, color=MUT_R)
    tb(slide, '→  Hardware attestation + sensor fusion + anomaly detection close the remaining gap between "math is correct" and "data is real".',
       ML, SH-Inches(0.6), SW-ML-Inches(0.45), Inches(0.4), size=9.5, italic=True, color=AMB_R)

def s06_claim(slide, img):
    eyebrow(slide, 'RideTrue — Claim Cycle')
    heading(slide, 'End-to-End Automated Claim Processing')
    tb(slide, 'Smart-contract enforced  ·  Oracle-verified  ·  USDC settled  ·  Transparent on-chain',
       ML, Inches(1.22), SW-ML-Inches(0.45), Inches(0.32), size=10.5, color=CYN_R)
    add_img(slide, img, ML, Inches(1.62), SW-ML-Inches(0.45), Inches(1.95))
    t = Inches(3.7)
    box_shape(slide, ML, t, SW-ML-Inches(0.45), Inches(3.1), '#0d1220', RED_H, 1.5)
    tb(slide, 'Fraud Guardrails', ML+Inches(0.2), t+Inches(0.15), SW-ML-Inches(0.7), Inches(0.35),
       size=12, bold=True, color=RED_R)
    guardrails = [
        'High claim amount + low rider reputation → automatic escalation to manual arbitration',
        'Repeated mechanic/rider pair anomalies → elevated dispute probability flag',
        'Proven fraud → stake slashing + permanent reputation penalty + retroactive fee clawback',
        'Appeals are time-bounded — no open-ended liability that threatens pool solvency',
    ]
    bxt = slide.shapes.add_textbox(ML+Inches(0.2), t+Inches(0.55), SW-ML-Inches(0.8), Inches(2.3))
    tf = bxt.text_frame; tf.word_wrap = True
    for j,g in enumerate(guardrails):
        p = tf.paragraphs[0] if j==0 else tf.add_paragraph()
        p.space_before = Pt(6)
        r = p.add_run(); r.text = f'• {g}'
        r.font.size = Pt(10.5); r.font.color.rgb = TXT_R

def s07_revenue(slide, img):
    eyebrow(slide, 'Economics — Revenue Model')
    heading(slide, 'Four Revenue Streams · Seoul Year 1')
    t = Inches(1.22)
    add_img(slide, img, Inches(0.15), t-Inches(0.1), Inches(6.1), Inches(5.9))
    x = Inches(6.5)
    rw = SW - x - Inches(0.45)
    rows = [
        ('VIN Registration',   '₩10,000 / bike',    '5,000 bikes',      '₩50 M'),
        ('Ownership Transfer', '₩15,000 / transfer','1,000 transfers',   '₩15 M'),
        ('Service Logging',    '₩5,000 / record',   '10,000 records',    '₩50 M'),
        ('Insurance Pool Fee', '5% of premiums',    '₩2 B pool',         '₩100 M'),
        ('Claims Processing',  '1.5% of claim',     '₩150 M payouts',    '₩2.25 M'),
        ('API / Data Access',  'Subscription',       'Pilot tier',        '₩50 M'),
    ]
    hdr_y = t
    for label, w, txt in [('Revenue Stream', Inches(2.4), 'stream'),
                           ('Unit Price', Inches(1.65), 'price'),
                           ('Y1 Revenue', Inches(1.2), 'rev')]:
        xoff = {'stream': 0, 'price': 2.45, 'rev': 4.15}[txt]
        tb(slide, label, x+Inches(xoff), hdr_y, Inches(w), Inches(0.3),
           size=9, bold=True, color=BLU_R)
    for i,(stream, price, vol, rev) in enumerate(rows):
        ry = hdr_y + Inches(0.35 + i*0.62)
        box_shape(slide, x, ry, rw, Inches(0.56), LIN_H)
        tb(slide, stream, x+Inches(0.12), ry+Inches(0.09), Inches(2.25), Inches(0.38), size=10, color=TXT_R)
        tb(slide, price,  x+Inches(2.45), ry+Inches(0.09), Inches(1.6),  Inches(0.38), size=9,  color=MUT_R)
        tb(slide, rev,    x+Inches(4.15), ry+Inches(0.09), Inches(1.15), Inches(0.38), size=10.5, bold=True, color=GRN_R)
    total_y = hdr_y + Inches(0.35 + len(rows)*0.62)
    box_shape(slide, x, total_y, rw, Inches(0.58), '#0d1f10', GRN_H, 1.5)
    tb(slide, 'TOTAL YEAR 1', x+Inches(0.12), total_y+Inches(0.1), Inches(2.5), Inches(0.38),
       size=11, bold=True, color=TXT_R)
    tb(slide, '₩ 227 M',     x+Inches(4.15), total_y+Inches(0.08), Inches(1.2), Inches(0.42),
       size=15, bold=True, color=GRN_R)
    tb(slide, 'Year 1 OpEx ≈ ₩1.8 B — scale and trust-building story, not a short-term profit pitch.',
       x, total_y+Inches(0.72), rw, Inches(0.38), size=9, italic=True, color=MUT_R)

def s08_token(slide, img):
    eyebrow(slide, 'Tokenomics — $MOTO')
    heading(slide, '$MOTO Token · 1 Billion Fixed Supply · No Post-TGE Inflation')
    t = Inches(1.22)
    add_img(slide, img, Inches(0.1), t-Inches(0.1), Inches(6.5), Inches(5.9))
    x = Inches(7.0); bw = SW-x-Inches(0.45)
    tb(slide, 'Token Utility', x, t, bw, Inches(0.35), size=13, bold=True, color=TXT_R)
    utils = [
        ('Staking & Slashing', GRN_H,
         'Lock $MOTO to gain write or verification rights. Bad behavior slashes your stake.'),
        ('Governance', BLU_H,
         'Vote on fees, reserve ratios, rebate formulas, treasury deployment, and integrations.'),
        ('Reputation Gating', AMB_H,
         'Tiered staking unlocks higher write limits and bigger fee shares for trusted participants.'),
        ('Fee Discounts', CYN_H,
         'Pay API and registration fees in $MOTO at a discount — fees burned or sent to treasury.'),
    ]
    for i,(title,col,desc) in enumerate(utils):
        uy = t + Inches(0.45 + i*1.28)
        box_shape(slide, x, uy, bw, Inches(1.12), LIN_H, col, 1.2)
        tb(slide, title, x+Inches(0.14), uy+Inches(0.1), bw, Inches(0.32), size=10.5, bold=True, color=_rgb(col))
        tb(slide, desc,  x+Inches(0.14), uy+Inches(0.46), bw-Inches(0.28), Inches(0.58), size=9.5, color=TXT_R)
    tb(slide, '⚠  All insurance returns and rider rebates are paid in USDC — not $MOTO. Deliberate: prevents reflexive demand loops.',
       x, t+Inches(5.6), bw, Inches(0.5), size=9, italic=True, color=AMB_R)

def s09_insurance(slide, img):
    eyebrow(slide, 'Economics — Insurance Pool')
    heading(slide, 'Layered Reserves · Where Every Premium Goes')
    t = Inches(1.22)
    add_img(slide, img, Inches(0.1), t-Inches(0.1), Inches(6.1), Inches(5.9))
    x = Inches(6.5); bw = SW-x-Inches(0.45)
    reserves = [
        ('Liquid Reserve  25%',      GRN_H, 'Always liquid. Never deployed to DeFi. Pays claims immediately.'),
        ('Catastrophe Reserve  20%', BLU_H, 'Separate fund for mass events — typhoon, flood, infrastructure outage.'),
        ('Safe-Rider Rebates  18%',  CYN_H, 'Returned quarterly to riders with zero or low claims.'),
        ('Expected Claims  12%',     RED_H, 'Pays approved claims. Buffer above modeled ~9% loss ratio.'),
        ('Reinsurance  10%',         AMB_H, 'External backstop — covers tail risk above pool capacity.'),
        ('Treasury + Fee  15%',      MUT_H, 'Protocol operational reserve and revenue.'),
    ]
    for i,(label,col,desc) in enumerate(reserves):
        ry = t + Inches(i*0.9)
        tb(slide, '●', x, ry+Inches(0.08), Inches(0.22), Inches(0.3), size=14, color=_rgb(col))
        tb(slide, label, x+Inches(0.25), ry+Inches(0.04), bw, Inches(0.3), size=10, bold=True, color=_rgb(col))
        tb(slide, desc,  x+Inches(0.25), ry+Inches(0.35), bw-Inches(0.1), Inches(0.45), size=9, color=MUT_R)

def s10_mechanic(slide, img):
    eyebrow(slide, 'Trust & Fraud Prevention')
    heading(slide, 'Three-Layer Defense That Makes Fraud Unprofitable')
    t = Inches(1.22)
    add_img(slide, img, ML, t, SW-ML-Inches(0.45), Inches(5.8))

def s11_pnl(slide, img):
    eyebrow(slide, 'Economics — Financial Model')
    heading(slide, '3-Year Stressed P&L · Break-Even in Year 3')
    tb(slide, 'Stress: 15% claim frequency · ₩2.5 M avg severity — credible delivery-rider scenario.',
       ML, Inches(1.22), SW-ML-Inches(0.45), Inches(0.32), size=10.5, italic=True, color=AMB_R)
    add_img(slide, img, ML, Inches(1.62), Inches(9.2), Inches(4.35))
    x = Inches(9.9); bw = SW-x-Inches(0.4)
    callouts = [
        ('₩ 4.5 B', 'Total seed + Series A\nfunding requirement', BLU_H),
        ('~₩ 2 B',  'Runway buffer beyond\nbreak-even', AMB_H),
        ('Year 3',  'First sustained\nprotocol profitability', GRN_H),
    ]
    for i,(kpi,label,col) in enumerate(callouts):
        cy = Inches(1.85 + i*1.65)
        box_shape(slide, x, cy, bw, Inches(1.4), LIN_H, col, 1.2)
        tb(slide, kpi,   x+Inches(0.14), cy+Inches(0.1),  bw, Inches(0.55), size=20, bold=True, color=_rgb(col))
        tb(slide, label, x+Inches(0.14), cy+Inches(0.68), bw, Inches(0.6),  size=9,  color=TXT_R)

def s12_lp(slide, img):
    eyebrow(slide, 'Economics — LP Capital Plan')
    heading(slide, 'Pool Solvency Requires Capital Beyond Rider Premiums')
    t = Inches(1.22)
    add_img(slide, img, ML, t, Inches(8.6), Inches(4.5))
    x = Inches(9.25); bw = SW-x-Inches(0.4)
    tb(slide, 'Solvency Ratio\nLP Capital / Annual Premiums', x, t, bw, Inches(0.65),
       size=10, bold=True, color=TXT_R)
    for i,(yr,ratio,col) in enumerate([('Year 1','1.39×',GRN_H),('Year 2','0.81×',AMB_H),('Year 3','0.68×',RED_H)]):
        ry = t + Inches(0.75 + i*1.15)
        box_shape(slide, x, ry, bw, Inches(1.0), LIN_H, col, 1.2)
        tb(slide, yr,    x+Inches(0.14), ry+Inches(0.06), bw, Inches(0.3),  size=10, color=MUT_R)
        tb(slide, ratio, x+Inches(0.14), ry+Inches(0.38), bw, Inches(0.45), size=20, bold=True, color=_rgb(col))
    tb(slide, '⚠  If only foundation seed materializes, solvency falls below 0.7× by Year 3 — insufficient for catastrophe coverage. LP partnerships must be confirmed before launch.',
       ML, t+Inches(4.7), SW-ML-Inches(0.45), Inches(0.55), size=9.5, italic=True, color=RED_R)

def s13_competitive(slide, img):
    eyebrow(slide, 'Competitive Analysis')
    heading(slide, 'Three Competitor Groups · One Intersection They Cannot Occupy')
    t = Inches(1.22)
    add_img(slide, img, ML, t, SW-ML-Inches(0.45), Inches(3.9))
    moat = [
        ('Vertical Data\nNetwork Effect', 'More bikes + riders + shops = more valuable dataset — compounds over time', BLU_H),
        ('Two-Wheelers Only', 'Deep vertical specialization vs. generic solutions from all three competitors', GRN_H),
        ('Identity + Behavior', 'Two protocols, one stack — each reinforces the other, none of them have both', CYN_H),
        ('Korea-First Wedge', 'Dense urban motorcycle market, replicable across Southeast Asia', AMB_H),
    ]
    mw = (SW-ML-Inches(0.45)) / 4
    for i,(title,desc,col) in enumerate(moat):
        x = ML + i*mw
        box_shape(slide, x, t+Inches(4.12), mw-Inches(0.12), Inches(1.9), LIN_H, col, 1.2)
        tb(slide, title, x+Inches(0.12), t+Inches(4.27), mw-Inches(0.25), Inches(0.55),
           size=10, bold=True, color=_rgb(col))
        tb(slide, desc,  x+Inches(0.12), t+Inches(4.85), mw-Inches(0.25), Inches(1.0),
           size=9, color=TXT_R)

def s14_security(slide, img):
    eyebrow(slide, 'Security & Risk Management')
    heading(slide, 'Security Is a Launch Gate — Not a Post-Launch Concern')
    tb(slide, 'The protocol will not accept user funds until both independent smart-contract audits are complete and all critical findings are fully resolved.',
       ML, Inches(1.22), SW-ML-Inches(0.45), Inches(0.38), size=11, bold=True, color=RED_R)
    add_img(slide, img, ML, Inches(1.68), SW-ML-Inches(0.45), Inches(3.8))
    tb(slide, 'Black Swan Plan — 9 safeguards: mandatory liquid reserve ≥25% · catastrophe fund 20% · reinsurance backstop · '
              'circuit breakers (claim spike / low reserve / disaster zone) · coverage caps · conservative DeFi yield (10–15% cap) · '
              'staged payouts (50% immediate / 50% over 3–6 months) · pro-rata distribution · temporary premium adjustment (+10–20%).',
       ML, Inches(5.65), SW-ML-Inches(0.45), Inches(1.0), size=9.5, italic=True, color=AMB_R)

def s15_closing(slide, eco_img):
    eyebrow(slide, 'Why MotoChain & RideTrue')
    heading(slide, 'The Intersection None of Them Can Occupy Alone', size=26)
    t = Inches(1.22)
    add_img(slide, eco_img, Inches(0.1), t, Inches(7.0), Inches(5.7))
    x = Inches(7.3); bw = SW-x-Inches(0.4)
    points = [
        ('Buildable today', BLU_H,
         'Base L2 + proven ZK tooling + DeFi composability — no new infrastructure needed.'),
        ('$190 B+ market', GRN_H,
         'Combined addressable market across global motorcycle resale and insurance.'),
        ('Defensible moat', CYN_H,
         'Vertical data network effects compound with every new participant on the protocol.'),
        ('User-owned', AMB_H,
         'Behavior scores and vehicle history follow the rider — not the platform. True portability.'),
    ]
    for i,(title,col,desc) in enumerate(points):
        py = t + Inches(i*1.35)
        box_shape(slide, x, py, bw, Inches(1.18), LIN_H, col, 1.2)
        tb(slide, title, x+Inches(0.16), py+Inches(0.1),  bw, Inches(0.35), size=13, bold=True, color=_rgb(col))
        tb(slide, desc,  x+Inches(0.16), py+Inches(0.5),  bw-Inches(0.2), Inches(0.58), size=9.5, color=TXT_R)
    tb(slide, 'Mireu Kim  ·  Enxhi Topalli  ·  CS495 DeFi  ·  May 2025',
       x, t+Inches(5.5), bw, Inches(0.35), size=9.5, color=MUT_R, align=PP_ALIGN.CENTER)

# ── BUILD ─────────────────────────────────────────────────────────────────

print('Generating charts …')
vin_img   = make_vin_flow()
zk_img    = make_zk_flow()
claim_img = make_claim_flow()
rev_img   = make_revenue_pie()
tok_img   = make_token_donut()
prem_img  = make_premium_pie()
mech_img  = make_mechanic_defense()
pnl_img   = make_pnl_bars()
lp_img    = make_lp_bars()
comp_img  = make_competitive()
sec_img   = make_security_timeline()
eco_img   = make_ecosystem()
print('  ✓ all charts done')

TOTAL = 15
prs = new_prs()

slides_plan = [
    ('Title',             s01_title),
    ('The Problem',       s02_problem),
    ('The Solution',      s03_solution),
    ('VIN Lifecycle',     lambda s: s04_vin(s, vin_img)),
    ('ZK Proof Flow',     lambda s: s05_zk(s, zk_img)),
    ('Claim Cycle',       lambda s: s06_claim(s, claim_img)),
    ('Revenue Model',     lambda s: s07_revenue(s, rev_img)),
    ('$MOTO Token',       lambda s: s08_token(s, tok_img)),
    ('Insurance Pool',    lambda s: s09_insurance(s, prem_img)),
    ('Mechanic Trust',    lambda s: s10_mechanic(s, mech_img)),
    ('3-Year P&L',        lambda s: s11_pnl(s, pnl_img)),
    ('LP Capital Plan',   lambda s: s12_lp(s, lp_img)),
    ('Competitive',       lambda s: s13_competitive(s, comp_img)),
    ('Security Roadmap',  lambda s: s14_security(s, sec_img)),
    ('Closing / Moat',    lambda s: s15_closing(s, eco_img)),
]

print('Building slides …')
for n, (name, builder) in enumerate(slides_plan, 1):
    print(f'  {n:2d}/{TOTAL}  {name}')
    sl = add_slide(prs)
    builder(sl)
    if n > 1:  # no page number on title
        pagenum(sl, n, TOTAL)

out = '0522_G1_Interim_Report.pptx'
prs.save(out)
print(f'\n✓ Saved → {out}  ({TOTAL} slides)')
