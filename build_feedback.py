#!/usr/bin/env python3
"""
MotoChain feedback build script.
Addresses interim-report presentation feedback across:
  - G1_MotoChain_RideTrue_INTERIM_REPORTpptx.pptx
  - MotoChain & Ride True Whitepaper.docx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

from docx import Document
from docx.shared import Pt as DPt, RGBColor as DRGB, Inches as DInches
from docx.oxml.ns import qn as dqn
from docx.oxml import OxmlElement
from docx.enum.text import WD_ALIGN_PARAGRAPH

# ─── Colour palette (matched to existing slide CSS) ─────────────────────────
BLUE   = RGBColor(0x4A, 0xA3, 0xFF)
CYAN   = RGBColor(0x27, 0xD8, 0xD8)
GREEN  = RGBColor(0x4E, 0xE5, 0x9A)
AMBER  = RGBColor(0xFF, 0xBF, 0x56)
RED    = RGBColor(0xFF, 0x6B, 0x6B)
WHITE  = RGBColor(0xED, 0xF4, 0xFF)
MUTED  = RGBColor(0xB8, 0xC8, 0xE0)
LABEL  = RGBColor(0x9C, 0xC6, 0xFF)
SUBVAL = RGBColor(0xC8, 0xDC, 0xFF)

# EMU shortcuts
I = Inches

# ─── PPTX helpers ────────────────────────────────────────────────────────────

def tb(slide, text, left, top, width, height,
       size=9.5, bold=False, color=MUTED,
       align=PP_ALIGN.LEFT, italic=False, space_before=0):
    """Add a text box with a single run."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def header(slide, eyebrow, title, subtitle=""):
    """Standard slide header (matches existing slide layout)."""
    tb(slide, eyebrow, I(0.70), I(0.40), I(8.00), I(0.32),
       size=11, bold=True, color=LABEL)
    tb(slide, title, I(0.70), I(0.74), I(12.00), I(0.72),
       size=36, bold=True, color=WHITE)
    if subtitle:
        tb(slide, subtitle, I(0.70), I(1.44), I(12.30), I(0.36),
           size=13, color=MUTED)


def footer(slide, note=""):
    """Standard footer bar (matches existing)."""
    if note:
        tb(slide, note, I(0.70), I(7.25), I(9.90), I(0.20),
           size=8.5, color=LABEL)
    tb(slide, "MotoChain & RideTrue", I(10.50), I(7.28), I(2.50), I(0.18),
       size=8.5, bold=True, color=LABEL, align=PP_ALIGN.RIGHT)


def sep(slide):
    """Thin separator line below header."""
    tb(slide, "", I(0.70), I(1.86), I(1.20), I(0.06), color=LABEL)


def clear_shapes(slide):
    """Remove all shapes from a slide (keeps mandatory XML skeleton)."""
    sp_tree = slide.shapes._spTree
    mandatory = {qn('p:nvGrpSpPr'), qn('p:grpSpPr')}
    for el in [c for c in sp_tree if c.tag not in mandatory]:
        sp_tree.remove(el)


def dup_slide(prs, idx):
    """Duplicate slide[idx] → append at end, return new slide."""
    src = prs.slides[idx]
    new = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    clear_shapes(new)
    for el in src.shapes._spTree:
        new.shapes._spTree.append(copy.deepcopy(el))
    bg = src._element.find(qn('p:bg'))
    if bg is not None:
        new._element.insert(2, copy.deepcopy(bg))
    return new


def blank_slide(prs):
    """Add a blank slide inheriting master background."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def move_slide(prs, old, new):
    """Move slide from index old to index new."""
    lst  = prs.slides._sldIdLst
    item = lst[old]
    lst.remove(item)
    lst.insert(new, item)


def find_shape_by_text(slide, substr):
    for shape in slide.shapes:
        if shape.has_text_frame and substr in shape.text_frame.text:
            return shape
    return None


def set_run_text(shape, old_substr, new_text):
    """Replace first run whose text contains old_substr."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old_substr in run.text:
                run.text = run.text.replace(old_substr, new_text)
                return True
    return False


# ════════════════════════════════════════════════════════════════════════════
#  PPTX CHANGES
# ════════════════════════════════════════════════════════════════════════════

def modify_slide7_biometrics(prs):
    """Add biometrics as layer 5 to the Telemetry Integrity slide."""
    s = prs.slides[6]

    # 1. Subtitle: "Four" → "Five"
    sub = find_shape_by_text(s, "Four defenses")
    if sub:
        set_run_text(sub, "Four defenses", "Five defenses")

    # 2. Replace the "Not absolute." note with biometrics + condensed note
    note = find_shape_by_text(s, "Not absolute.")
    if note:
        # Move up 0.22" and expand height
        note.top    = I(6.86)
        note.height = I(0.50)
        tf = note.text_frame
        tf.word_wrap = True
        # Clear existing paragraphs
        for para in tf.paragraphs:
            for run in para.runs:
                run.text = ""
        # First para: biometrics
        p0 = tf.paragraphs[0]
        r0 = p0.runs[0] if p0.runs else p0.add_run()
        r0.text = (
            "5  ·  BIOMETRIC BINDING  —  tie each session to the registered rider "
            "(Face ID, fingerprint, or voice). Prevents score farming: a careful "
            "rider cannot cover for a riskier one."
        )
        r0.font.size      = Pt(9)
        r0.font.bold      = False
        r0.font.color.rgb = MUTED
        # Second para: original note (condensed)
        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = "No system is absolute — goal: detection × penalty > expected payoff."
        r1.font.size      = Pt(8.5)
        r1.font.color.rgb = LABEL


def modify_slide25_permissionless(prs):
    """Add permissionless onboarding note to Mechanic Trust slide."""
    s = prs.slides[24]

    # Find footer note shape and update text
    note = find_shape_by_text(s, "Each layer covers")
    if note:
        tf = note.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if "Each layer covers" in run.text:
                    run.text = (
                        "Each layer covers what the others miss — none works alone, "
                        "all three make fraud unprofitable.  ·  Registration is open: "
                        "anyone can stake $MOTO and begin logging. No business licence "
                        "required — your stake is your credential."
                    )
                    run.font.size = Pt(8.5)
                    return


# ─── NEW SLIDE: Why Two Protocols ────────────────────────────────────────────

def create_why_two_protocols(prs):
    s = blank_slide(prs)

    header(s,
           "II.2  ·  PROTOCOL DESIGN",
           "ONE CHAIN · TWO CONTRACTS · TWO TYPES OF TRUTH",
           "VIN follows the bike. Riding score follows the rider. That separation is the design — not an accident.")
    sep(s)

    # ── Left panel: MotoChain VIN ────────────────────────────────────────────
    tb(s, "MOTOCHAIN  ·  VIN IDENTITY TOKEN",
       I(0.50), I(2.05), I(5.90), I(0.30),
       size=10, bold=True, color=CYAN)

    tb(s, "TRANSFERRABLE",
       I(0.50), I(2.38), I(5.90), I(0.50),
       size=26, bold=True, color=CYAN)

    tb(s, "The motorcycle's permanent identity token. Minted once, travels with the asset.",
       I(0.50), I(2.92), I(5.90), I(0.40),
       size=11, color=WHITE)

    items_vin = [
        "Service records (co-signed by mechanic + owner)",
        "Full ownership chain  —  title is the wallet",
        "Accident log & mileage history",
        "Fraud-detection flags and escrow outcomes",
    ]
    y = 3.38
    for item in items_vin:
        tb(s, f"  →  {item}", I(0.50), I(y), I(5.90), I(0.28),
           size=10.5, color=MUTED)
        y += 0.30

    tb(s, (
        "When you sell the bike  →  the full verified history transfers "
        "to the new owner. The buyer gets everything the bike has ever "
        "accumulated."
       ),
       I(0.50), I(4.62), I(5.90), I(0.70),
       size=10.5, color=SUBVAL)

    tb(s, "WHY SEPARATE?",
       I(0.50), I(5.42), I(5.90), I(0.28),
       size=9, bold=True, color=CYAN)
    tb(s, (
        "The new owner needs the bike's full history — but they build "
        "their own riding reputation from scratch on RideTrue."
       ),
       I(0.50), I(5.72), I(5.90), I(0.50),
       size=10, color=MUTED)

    # ── Divider ──────────────────────────────────────────────────────────────
    tb(s, "⟵  bike  |  rider  ⟶",
       I(6.17), I(3.80), I(1.05), I(0.80),
       size=8.5, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

    # ── Right panel: RideTrue Score ──────────────────────────────────────────
    tb(s, "RIDETRUE  ·  RIDER SAFETY SCORE",
       I(7.40), I(2.05), I(5.60), I(0.30),
       size=10, bold=True, color=GREEN)

    tb(s, "SOULBOUND",
       I(7.40), I(2.38), I(5.60), I(0.50),
       size=26, bold=True, color=GREEN)

    tb(s, "The rider's personal reputation. Built by real rides, attested by ZK proof. Cannot be sold, transferred, or delegated.",
       I(7.40), I(2.92), I(5.60), I(0.40),
       size=11, color=WHITE)

    items_rider = [
        "Riding behaviour score (speed, braking, cornering)",
        "ZK-attested proof — no raw data exposed",
        "Claims history & penalty flags",
        "Biometric session binding (Face ID / fingerprint)",
    ]
    y = 3.38
    for item in items_rider:
        tb(s, f"  →  {item}", I(7.40), I(y), I(5.60), I(0.28),
           size=10.5, color=MUTED)
        y += 0.30

    tb(s, (
        "When you sell the bike  →  you keep your score. New bike, "
        "same reputation. Any insurer on the protocol sees your "
        "portable, verified riding history."
       ),
       I(7.40), I(4.62), I(5.60), I(0.70),
       size=10.5, color=SUBVAL)

    tb(s, "WHY SEPARATE?",
       I(7.40), I(5.42), I(5.60), I(0.28),
       size=9, bold=True, color=GREEN)
    tb(s, (
        "Rider reputation is personal — transferring it with the bike "
        "would let bad riders hide behind a 'clean' vehicle history."
       ),
       I(7.40), I(5.72), I(5.60), I(0.50),
       size=10, color=MUTED)

    footer(s, "Result: buyers get verified bike history. Riders keep a portable score that earns fair pricing from any insurer on the protocol.")
    return s


# ─── NEW SLIDE: Participant Interaction Flow ─────────────────────────────────

def create_participant_flow(prs):
    s = blank_slide(prs)

    header(s,
           "II.4  ·  ECOSYSTEM DYNAMICS",
           "HOW PARTICIPANTS INTERACT",
           "Every action creates value that flows to others — that is what makes this a protocol, not a product.")
    sep(s)

    # Columns: Riders | Mechanics | Dealers | LPs | Insurers/Lenders
    cols = [
        ("RIDERS",             CYAN,  0.35, "Premium deposits\n+ ride data", "Claims paid\n+ safe-rider rebates\n+ MOTO stake rewards"),
        ("MECHANICS",          GREEN, 2.90, "Service entries\n+ MOTO stake", "60% fee share\n+ MOTO rewards\nfor quality streaks"),
        ("DEALERS",            AMBER, 5.45, "VIN registration\n+ transfer fees", "Verified inventory\n+ MOTO early-adopter\nrewards"),
        ("LPs",                BLUE,  7.95, "USDC into\ninsurance pool", "Yield from DeFi\n+ fee share"),
        ("INSURERS\n& LENDERS", RED,  10.35, "Data subscription\nAPI fees", "Verified VIN history\n+ privacy-preserving\nrider risk scores"),
    ]

    for name, color, x, gives, gets in cols:
        # Actor label
        tb(s, name, I(x), I(2.08), I(2.55), I(0.42),
           size=10, bold=True, color=color, align=PP_ALIGN.CENTER)
        # Arrow down
        tb(s, "↓", I(x + 1.0), I(2.55), I(0.55), I(0.26),
           size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
        # What they give
        tb(s, gives, I(x), I(2.82), I(2.55), I(0.70),
           size=9.5, color=color, align=PP_ALIGN.CENTER)

    # Central protocol box
    tb(s, "MOTOCHAIN\n& RIDETRUE\nPROTOCOL",
       I(5.42), I(3.68), I(2.50), I(0.90),
       size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for name, color, x, gives, gets in cols:
        # What they get back
        tb(s, gets, I(x), I(4.72), I(2.55), I(0.80),
           size=9.5, color=color, align=PP_ALIGN.CENTER)
        # Arrow up
        tb(s, "↑", I(x + 1.0), I(4.54), I(0.55), I(0.26),
           size=14, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Network-effect note
    tb(s, (
        "Data flywheel:  more bikes + riders + mechanics  →  richer, more "
        "accurate data  →  better risk pricing  →  lower premiums  →  more "
        "participants  →  stronger dataset.  Each new entry makes every "
        "other entry more valuable."
       ),
       I(0.50), I(5.68), I(12.30), I(0.58),
       size=10.5, color=MUTED, align=PP_ALIGN.CENTER)

    footer(s)
    return s


# ─── NEW SLIDE: Dealer Flywheel / Second-Hand Market ────────────────────────

def create_dealer_flywheel(prs):
    s = blank_slide(prs)

    header(s,
           "II.4  ·  DEALER STRATEGY",
           "THE SECOND-HAND DEALER OPPORTUNITY",
           "Objection: \"My market is already doing well — why register older motorcycles?\"")
    sep(s)

    # ── Left: The economic answer ─────────────────────────────────────────────
    tb(s, "THE ANSWER: VERIFIED BIKES SELL FOR MORE",
       I(0.50), I(2.05), I(6.00), I(0.30),
       size=10, bold=True, color=CYAN)

    rows = [
        ("One-time VIN registration fee",   "₩10,000"),
        ("Typical price premium on verified bike", "₩200K – ₩500K+"),
        ("Net return on that fee",          "20× – 50×"),
    ]
    y = 2.42
    for label, val in rows:
        tb(s, label, I(0.50), I(y), I(4.40), I(0.32), size=11, color=WHITE)
        tb(s, val,   I(5.00), I(y), I(1.40), I(0.32),
           size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
        y += 0.35

    tb(s, (
        "For older bikes specifically: buyers face the greatest uncertainty "
        "in precisely that segment, which means verified condition is most "
        "valuable exactly where dealers need the pricing edge."
       ),
       I(0.50), I(3.52), I(6.00), I(0.60),
       size=10.5, color=MUTED)

    tb(s, "RETROACTIVE CONDITION ANCHORING",
       I(0.50), I(4.22), I(6.00), I(0.28),
       size=10, bold=True, color=CYAN)

    tb(s, (
        "Even without a full service history, a current-condition "
        "certification has real market value. MotoChain documents what "
        "the bike is today. Every future service record builds on that "
        "anchor, compounding the asset's verified value over time."
       ),
       I(0.50), I(4.54), I(6.00), I(0.80),
       size=10.5, color=WHITE)

    # ── Right: Incentives + flywheel ─────────────────────────────────────────
    tb(s, "EARLY-ADOPTER INCENTIVES",
       I(6.70), I(2.05), I(6.30), I(0.30),
       size=10, bold=True, color=AMBER)

    incentives = [
        ("MOTO rewards",
         "Dealers who register ≥50 older bikes in the launch cohort earn "
         "$MOTO rewards equivalent to 12 months of fee rebates."),
        ("Verified Dealer badge",
         "Displayed on partner marketplaces — visible differentiation that "
         "grows with every registration, compounding trust over time."),
        ("Tier-based fee reductions",
         "Higher verified inventory unlocks lower protocol fees on future "
         "ownership transfers — a direct financial incentive to keep growing."),
    ]
    y = 2.42
    for title, body in incentives:
        tb(s, title, I(6.70), I(y), I(6.20), I(0.22),
           size=10, bold=True, color=SUBVAL)
        tb(s, body, I(6.70), I(y + 0.24), I(6.20), I(0.38),
           size=10, color=MUTED)
        y += 0.72

    tb(s, "THE FLYWHEEL",
       I(6.70), I(4.70), I(6.30), I(0.28),
       size=10, bold=True, color=GREEN)

    tb(s, (
        "Dealer registers old bike  →  buyer sees verified condition  →  "
        "buyer pays premium  →  dealer earns more per sale  →  dealer "
        "registers more bikes  →  protocol data grows  →  pricing improves "
        "for all participants."
       ),
       I(6.70), I(5.02), I(6.20), I(0.70),
       size=10.5, color=WHITE)

    footer(s,
           "Core reframe: MotoChain is not extra paperwork — "
           "it is the tool that turns old bikes into premium-priced inventory.")
    return s


# ─── NEW SLIDE: $MOTO Utility & Community Incentives ────────────────────────

def create_moto_utility(prs):
    s = blank_slide(prs)

    header(s,
           "II.5  ·  TOKEN UTILITY",
           "$MOTO — EARN, NOT JUST HOLD",
           "Like Uniswap LP fees: put $MOTO to work in the protocol and earn from real activity — not token emissions.")
    sep(s)

    # ── Four earner groups ───────────────────────────────────────────────────
    groups = [
        (CYAN,  "RIDERS",
         "Maintain a high safety score + zero fraudulent claims",
         "Stablecoin safe-rider rebates.\nMOTO stake unlocks higher rebate tiers and priority dispute resolution."),
        (GREEN, "MECHANICS",
         "Stake $MOTO → log verified service records",
         "60% of ₩5,000 per entry (immediate fee share).\n+MOTO rewards for quality-record streaks. Bad records slash the stake."),
        (AMBER, "DEALERS",
         "Register bikes — especially older legacy inventory",
         "MOTO from the 30% community ecosystem bucket.\nHigher verified count = lower protocol fees on future transfers."),
        (BLUE,  "VALIDATORS & ORACLES",
         "Stake $MOTO to validate claims and data feeds",
         "Share of claim-processing fees on every approved claim.\nSlashed for bad data. Rewarded for accurate, timely consensus."),
    ]

    x = 0.50
    col_w = 3.00
    for color, title, action, earn in groups:
        tb(s, title,  I(x), I(2.05), I(col_w), I(0.28), size=10, bold=True, color=color)
        tb(s, "Action:", I(x), I(2.38), I(col_w), I(0.20), size=8.5, bold=True, color=LABEL)
        tb(s, action,  I(x), I(2.60), I(col_w), I(0.46), size=9.5, color=WHITE)
        tb(s, "Earn:", I(x), I(3.10), I(col_w), I(0.20), size=8.5, bold=True, color=color)
        tb(s, earn,   I(x), I(3.32), I(col_w), I(0.70), size=9.5, color=MUTED)
        x += col_w + 0.10

    # ── Divider ──────────────────────────────────────────────────────────────
    tb(s, "", I(0.50), I(4.15), I(12.30), I(0.04), color=LABEL)

    # ── MOTO as fees (left) ──────────────────────────────────────────────────
    tb(s, "MOTO AS FEES",
       I(0.50), I(4.28), I(5.90), I(0.28),
       size=10, bold=True, color=AMBER)

    tb(s, (
        "VIN registration, ownership transfer, and API access fees can all "
        "be paid in $MOTO at a 20% discount versus stablecoin. The protocol "
        "burns a share of received $MOTO and routes the rest to the treasury "
        "— reducing stablecoin friction for early participants and creating a "
        "deflationary sink for the token."
       ),
       I(0.50), I(4.60), I(5.90), I(1.00),
       size=10.5, color=WHITE)

    # ── Permissionless governance (right) ────────────────────────────────────
    tb(s, "GOVERNANCE = PERMISSIONLESS FEES",
       I(6.60), I(4.28), I(6.10), I(0.28),
       size=10, bold=True, color=CYAN)

    tb(s, (
        "Fee rates are not set by founders — they are voted by $MOTO holders. "
        "Any holder can submit a governance proposal to adjust fees, reserve "
        "ratios, or reward formulas. Governance enforces permissionless "
        "parameter control: the protocol serves its community, not a central "
        "admin. Mechanics need no business licence — stake is the credential."
       ),
       I(6.60), I(4.60), I(6.10), I(1.00),
       size=10.5, color=WHITE)

    footer(s,
           "Your stake is your entry credential. Earn from the activity you "
           "create — mechanics, dealers, validators, and riders all have a "
           "direct economic reason to contribute quality data.")
    return s


# ════════════════════════════════════════════════════════════════════════════
#  DOCX CHANGES
# ════════════════════════════════════════════════════════════════════════════

def find_heading_idx(doc, heading_text):
    """Return index of first paragraph whose text contains heading_text."""
    for i, p in enumerate(doc.paragraphs):
        if heading_text in p.text:
            return i
    return -1


def insert_paragraph_after(doc, ref_idx, text, style="Normal", bold=False,
                            color=None, size=None, space_before=None):
    """Insert a new paragraph after ref_idx in doc.paragraphs."""
    ref_para = doc.paragraphs[ref_idx]
    new_para = OxmlElement('w:p')
    ref_para._element.addnext(new_para)
    # Find the newly added paragraph object
    for i, p in enumerate(doc.paragraphs):
        if p._element is new_para:
            idx = i
            break
    else:
        return None
    p = doc.paragraphs[idx]
    p.style = doc.styles[style]
    if space_before is not None:
        pPr = p._element.get_or_add_pPr()
        spacing = OxmlElement('w:spacing')
        spacing.set(dqn('w:before'), str(int(space_before * 20)))
        pPr.append(spacing)
    run = p.add_run(text)
    run.bold = bold
    if color:
        run.font.color.rgb = DRGB(*[int(color[i:i+2], 16) for i in (0, 2, 4)])
    if size:
        run.font.size = DPt(size)
    return idx


def add_heading_after(doc, ref_idx, text, level=3):
    """Insert a heading paragraph after ref_idx."""
    ref_para = doc.paragraphs[ref_idx]
    new_para = OxmlElement('w:p')
    ref_para._element.addnext(new_para)
    for i, p in enumerate(doc.paragraphs):
        if p._element is new_para:
            p.style = doc.styles[f'Heading {level}']
            p.add_run(text)
            return i
    return -1


def update_docx(path):
    doc = Document(path)

    # ── 1. Section 2.1 Telematics: add biometrics ────────────────────────────
    idx = find_heading_idx(doc, "2.1 Telematics Integrity")
    if idx >= 0:
        # Find the paragraph describing the 4 layers (contains "Sampling")
        search_idx = idx
        for j in range(idx, min(idx + 30, len(doc.paragraphs))):
            if "Sampling" in doc.paragraphs[j].text and "challenge" in doc.paragraphs[j].text.lower():
                search_idx = j
                break
        bio_text = (
            "A fifth recommended defense layer is biometric binding: "
            "Face ID, fingerprint, or voice confirmation ties each riding session "
            "to the registered individual rather than merely the device. This "
            "addresses a specific attack vector — score farming — where a careful "
            "rider operates a device on behalf of a riskier one to earn undeserved "
            "premium discounts. Biometric confirmation at session start makes this "
            "economically irrational by linking each score update to a verified "
            "human identity, not just a hardware token. Implementation should use "
            "the platform-native biometric APIs (iOS Face ID / Touch ID, Android "
            "BiometricPrompt) that are already hardware-attested, adding minimal "
            "friction while substantially closing the social-engineering gap."
        )
        insert_paragraph_after(doc, search_idx, bio_text)

    # ── 2. After architecture intro: VIN vs Score design rationale ───────────
    idx = find_heading_idx(doc, "2. Technical Architecture")
    if idx >= 0:
        # Find the paragraph mentioning "MotoChain lies a set of smart contracts"
        for j in range(idx, min(idx + 10, len(doc.paragraphs))):
            if "VIN-linked identity token" in doc.paragraphs[j].text:
                ins = j
                break
        else:
            ins = idx + 2

        vin_vs_score = (
            "A core design decision underpins the two-protocol architecture: the VIN "
            "identity token is transferrable, while the rider safety score is "
            "soulbound. These two data types represent fundamentally different "
            "assets. The VIN NFT follows the motorcycle: when ownership changes, the "
            "full verified history — service records, accident log, mileage chain — "
            "transfers to the new owner, because the buyer's due-diligence need "
            "requires access to that history. The rider score, by contrast, follows "
            "the person: it represents that individual's demonstrated behaviour "
            "behind a specific type of vehicle, accumulated over time, and cannot be "
            "sold or transferred without destroying its informational value. "
            "Conflating the two would allow a high-scoring rider to sell their "
            "reputation along with the bike, undermining the actuarial integrity of "
            "the insurance pool. The two-contract architecture enforces this "
            "separation at the protocol level, not merely as a policy."
        )
        add_heading_after(doc, ins - 1, "Why Two Contracts: VIN Token vs. Rider Score", level=3)
        idx2 = find_heading_idx(doc, "Why Two Contracts: VIN Token vs. Rider Score")
        if idx2 >= 0:
            insert_paragraph_after(doc, idx2, vin_vs_score)

    # ── 3. Section 4 Target Audience: dealer + old bikes ────────────────────
    idx = find_heading_idx(doc, "4. Target Audience")
    if idx >= 0:
        # Find the paragraph about "dealer" strategy
        for j in range(idx, min(idx + 25, len(doc.paragraphs))):
            if "dealer" in doc.paragraphs[j].text.lower() and "Seoul" in doc.paragraphs[j].text:
                ins = j
                break
        else:
            ins = idx + 4

        dealer_section_heading = "Dealer Incentive Strategy: The Second-Hand Market"
        add_heading_after(doc, ins, dealer_section_heading, level=3)
        h_idx = find_heading_idx(doc, dealer_section_heading)
        if h_idx >= 0:
            dealer_text = (
                "A critical objection from second-hand dealers must be answered "
                "directly: 'My resale market is already functioning. Why do extra "
                "work to register older motorcycles?' The protocol's response is "
                "economic, not ideological. Verified motorcycles command a price "
                "premium that materially exceeds the registration fee. A ₩10,000 "
                "one-time fee that enables a ₩200,000 to ₩500,000 price premium "
                "represents a 20× to 50× return on that cost. Dealers who register "
                "inventory — including older bikes — sell faster, at higher prices, "
                "and with fewer post-sale disputes. The incentive structure is "
                "designed to make this concrete. Dealers in the launch cohort who "
                "register a minimum of fifty older motorcycles earn $MOTO rewards "
                "equivalent to twelve months of fee rebates. A Verified Dealer badge "
                "is displayed on partner marketplaces, compounding the trust "
                "advantage with every additional registration. Higher verified "
                "inventory volumes unlock lower protocol fees on future ownership "
                "transfers, creating a permanent cost advantage for active dealers. "
                "For older motorcycles specifically, the value proposition is "
                "strongest precisely because buyer uncertainty is highest in that "
                "segment. The protocol addresses this through retroactive condition "
                "anchoring: even without a complete service history, a current "
                "condition certification documents the bike's verified state at a "
                "known point in time. Every subsequent service record builds on that "
                "anchor, compounding the asset's verifiable value. The dealer's "
                "market improves not despite older inventory, but because of it."
            )
            insert_paragraph_after(doc, h_idx, dealer_text)

    # ── 4. Section 5.5 Token Design: expand utility + MOTO as fees ──────────
    idx = find_heading_idx(doc, "5.5 Token Design")
    if idx >= 0:
        # Find the Token Utility paragraph
        for j in range(idx, min(idx + 30, len(doc.paragraphs))):
            if "Token Utility" in doc.paragraphs[j].text:
                ins = j
                break
        else:
            ins = idx + 5

        utility_heading = "5.5a Community Earning Mechanics and Permissionless Fees"
        add_heading_after(doc, ins, utility_heading, level=3)
        u_idx = find_heading_idx(doc, utility_heading)
        if u_idx >= 0:
            utility_text = (
                "The token design follows the Uniswap LP model: participants who "
                "contribute real value to the protocol earn from the economic "
                "activity they generate, not from speculative emissions. Four "
                "earner groups are defined. Mechanics stake $MOTO to write service "
                "records and earn sixty percent of each ₩5,000 entry fee as "
                "immediate income, supplemented by $MOTO rewards that scale with "
                "quality-record streaks. Fraudulent entries trigger stake slashing, "
                "creating a direct financial disincentive against false records. "
                "Dealers who register motorcycles — particularly older legacy "
                "inventory during the launch cohort period — earn $MOTO from the "
                "thirty-percent community ecosystem allocation. Higher verified "
                "inventory volumes unlock reduced transfer fees on future "
                "transactions, creating a permanent structural advantage. "
                "Validators and oracle operators stake $MOTO to confirm claims "
                "and data feeds and earn a share of claim-processing fees on each "
                "approved payout. Inaccurate feeds are slashed, accurate consensus "
                "is rewarded. Riders who maintain high safety scores and file no "
                "fraudulent claims receive stablecoin safe-rider rebates; $MOTO "
                "stake unlocks higher rebate tiers and priority dispute resolution. "
                "\n\n"
                "Fee payments in $MOTO are explicitly supported across all "
                "protocol interactions. VIN registration fees, ownership transfer "
                "fees, and B2B API access fees may each be settled in $MOTO at a "
                "twenty percent discount versus stablecoin. The protocol burns a "
                "portion of received $MOTO and routes the remainder to the "
                "treasury DAO, creating a deflationary token sink tied directly "
                "to real protocol usage. This reduces stablecoin friction for "
                "early participants and gives $MOTO holders a meaningful reason "
                "to hold rather than immediately liquidate."
                "\n\n"
                "Fee rates across the protocol are not fixed by the founding team. "
                "They are governed by $MOTO token holders through on-chain "
                "proposals. Any holder may submit a governance proposal to adjust "
                "fees, reserve ratios, reward formulas, or supported integrations. "
                "This permissionless parameter control is what distinguishes the "
                "protocol from a traditional financial intermediary: no central "
                "admin can unilaterally change the economics. The founding team "
                "retains no special governance authority beyond their token "
                "allocation, which vests on the same schedule as all other "
                "insiders."
            )
            insert_paragraph_after(doc, u_idx, utility_text)

    # ── 5. Section 5.6 Mechanic Trust: permissionless onboarding ────────────
    idx = find_heading_idx(doc, "5.6 Mechanic Trust")
    if idx >= 0:
        # Insert at beginning of this section
        perm_heading = "5.6a Permissionless Mechanic Registration"
        add_heading_after(doc, idx, perm_heading, level=3)
        p_idx = find_heading_idx(doc, perm_heading)
        if p_idx >= 0:
            perm_text = (
                "Mechanic registration on MotoChain is intentionally permissionless. "
                "No business licence, professional certification, or government "
                "registration is required to begin participating as a mechanic on "
                "the protocol. The entry credential is economic, not institutional: "
                "any participant who stakes $MOTO gains write access to service "
                "records at the Basic tier. This design choice is deliberate. "
                "Requiring licence verification would introduce a centralised "
                "gating mechanism that contradicts the protocol's core premise, "
                "create geographic and regulatory fragmentation that limits "
                "adoption, and exclude informal mechanics who represent a "
                "substantial share of the real-world motorcycle service market. "
                "\n\n"
                "Trust is enforced economically rather than institutionally. A "
                "mechanic with staked $MOTO has financial skin in the game: every "
                "service entry places a portion of stake at risk during the escrow "
                "window. Proven fraud triggers slashing. The bounty mechanism "
                "incentivises community members to surface fraudulent records. "
                "Tiered staking allows mechanics to voluntarily submit additional "
                "credentials — including business registration certificates, trade "
                "qualifications, or third-party reputation signals — to unlock "
                "higher tiers with greater write limits and fee shares. But the "
                "baseline tier is open to anyone, anywhere, who can stake. The "
                "protocol does not need to verify who you are; it verifies that "
                "you have something to lose."
            )
            insert_paragraph_after(doc, p_idx, perm_text)

    # ── Save ─────────────────────────────────────────────────────────────────
    out = path.replace(".docx", "_v2.docx")
    doc.save(out)
    print(f"  Whitepaper saved → {out}")
    return out


# ════════════════════════════════════════════════════════════════════════════
#  MAIN
# ════════════════════════════════════════════════════════════════════════════

def main():
    pptx_path = "G1_MotoChain_RideTrue_INTERIM_REPORTpptx.pptx"
    docx_path = "MotoChain & Ride True Whitepaper.docx"

    print("Loading presentation…")
    prs = Presentation(pptx_path)

    # ── Modify existing slides ───────────────────────────────────────────────
    print("  Modifying slide 7 (biometrics)…")
    modify_slide7_biometrics(prs)

    print("  Modifying slide 25 (permissionless mechanic note)…")
    modify_slide25_permissionless(prs)

    # ── Create new slides (appended at end) ─────────────────────────────────
    print("  Creating 'Why Two Protocols' slide…")
    create_why_two_protocols(prs)   # idx 25

    print("  Creating 'Participant Interaction Flow' slide…")
    create_participant_flow(prs)    # idx 26

    print("  Creating 'Dealer Flywheel' slide…")
    create_dealer_flywheel(prs)     # idx 27

    print("  Creating '$MOTO Utility' slide…")
    create_moto_utility(prs)        # idx 28

    # ── Reorder slides to correct positions ─────────────────────────────────
    # Current order: [0..24, A=25, B=26, C=27, D=28]
    # Target insertions:
    #   A after slide 5 (original idx 4) → position 5
    #   B after slide 11 "How It Works" (original idx 10) → becomes 12 after A
    #   C directly after B → position 13
    #   D after slide 18 "Token Allocation" (original idx 17) → becomes 21
    #                      after A,B,C inserted before it

    print("  Reordering slides…")

    # Step 1: Move A (currently at 25) → position 5
    move_slide(prs, 25, 5)
    # State: [0,1,2,3,4,A,5,6,7,8,9,10,11,12,13,14,15,16,17,18,19,20,21,22,23,24,B,C,D]

    # Step 2: Move B (currently at 26) → position 12
    # "How It Works" (orig idx 10) is now at idx 11 after A; insert after = 12
    move_slide(prs, 26, 12)
    # State: [0,1,2,3,4,A,5,6,7,8,9,10,11,B,12,13,14,15,16,17,18,19,20,21,22,23,24,C,D]

    # Step 3: Move C (currently at 27) → position 13 (directly after B)
    move_slide(prs, 27, 13)
    # State: [0,1,2,3,4,A,5,6,7,8,9,10,11,B,C,12,13,14,15,16,17,18,19,20,21,22,23,24,D]

    # Step 4: Move D (currently at 28) → position 22
    # Token Allocation (orig idx 17) is now at idx 20; insert after = 21
    move_slide(prs, 28, 21)

    print(f"  Total slides: {len(prs.slides)}")

    # ── Save PPTX ────────────────────────────────────────────────────────────
    out_pptx = pptx_path.replace(".pptx", "_v2.pptx")
    prs.save(out_pptx)
    print(f"  Presentation saved → {out_pptx}")

    # ── Update whitepaper ────────────────────────────────────────────────────
    print("Loading whitepaper…")
    update_docx(docx_path)

    print("\nDone.")


if __name__ == "__main__":
    main()
