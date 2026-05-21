#!/usr/bin/env python3
"""
MotoChain & RideTrue — Final Presentation
Follows whitepaper sequence exactly. Every slide has a branded header.
20 slides with matplotlib charts, flowcharts, and clean dark theme.
"""

import io, numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
from matplotlib.lines import Line2D

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────
BG_H  = '#060c16'; BG2_H = '#0d1624'; CARD_H = '#111c2e'
BLU_H = '#4aa3ff'; GRN_H = '#4ee59a'; AMB_H  = '#ffbf56'
CYN_H = '#27d8d8'; RED_H = '#ff6b6b'; PUR_H  = '#a78bfa'
TXT_H = '#e8eaf0'; MUT_H = '#8892a0'; LIN_H  = '#1e2a3a'

def _c(h):
    h=h.lstrip('#')
    return RGBColor(int(h[:2],16),int(h[2:4],16),int(h[4:],16))

BG_R=_c(BG_H); BG2_R=_c(BG2_H); CARD_R=_c(CARD_H)
BLU_R=_c(BLU_H); GRN_R=_c(GRN_H); AMB_R=_c(AMB_H)
CYN_R=_c(CYN_H); RED_R=_c(RED_H); TXT_R=_c(TXT_H)
MUT_R=_c(MUT_H); LIN_R=_c(LIN_H)

SW=Inches(13.333); SH=Inches(7.5); DPI=150
# Header bar metrics
HDR_H  = Inches(0.42)          # height of top header bar
EYEBROW_Y = HDR_H + Inches(0.1)
TITLE_Y   = HDR_H + Inches(0.38)
BODY_Y    = HDR_H + Inches(1.05)
ML = Inches(0.45); MR = Inches(0.45)
CW = SW - ML - MR              # content width

plt.rcParams.update({
    'figure.facecolor':BG_H,'axes.facecolor':BG2_H,'savefig.facecolor':BG_H,
    'text.color':TXT_H,'axes.labelcolor':TXT_H,'xtick.color':MUT_H,
    'ytick.color':MUT_H,'axes.edgecolor':LIN_H,'grid.color':LIN_H,
    'font.family':'sans-serif',
})

# ── matplotlib helpers ─────────────────────────────────────────────────────
def f2s(fig):
    buf=io.BytesIO()
    fig.savefig(buf,format='png',dpi=DPI,bbox_inches='tight',
                facecolor=fig.get_facecolor())
    plt.close(fig); buf.seek(0); return buf

def flow_chart(steps, fig_w=12.6, fig_h=2.5, bw=0.155, gap=0.035):
    """Generic left-to-right flow with numbered boxes."""
    fig,ax=plt.subplots(figsize=(fig_w,fig_h),facecolor=BG_H)
    ax.set_facecolor(BG_H); ax.axis('off')
    n=len(steps); total=n*bw+(n-1)*gap; sx=(1-total)/2
    for i,(num,label,body,col) in enumerate(steps):
        x=sx+i*(bw+gap)
        ax.add_patch(FancyBboxPatch((x,0.08),bw,0.84,boxstyle="round,pad=0.02",
            fc=col+'1e',ec=col,lw=2,transform=ax.transAxes,clip_on=False))
        ax.text(x+bw/2,0.79,num, transform=ax.transAxes,ha='center',va='center',
                fontsize=8,fontweight='bold',color=col)
        ax.text(x+bw/2,0.60,label,transform=ax.transAxes,ha='center',va='center',
                fontsize=8.5,fontweight='bold',color=TXT_H,linespacing=1.3)
        ax.text(x+bw/2,0.28,body, transform=ax.transAxes,ha='center',va='center',
                fontsize=7,color=MUT_H,linespacing=1.4)
        if i<n-1:
            ax.annotate('',xy=(x+bw+gap-0.006,0.5),xytext=(x+bw+0.005,0.5),
                xycoords='axes fraction',textcoords='axes fraction',
                arrowprops=dict(arrowstyle='->',color=MUT_H,lw=1.8))
    return f2s(fig)

# ── Charts ─────────────────────────────────────────────────────────────────
def ch_vin():
    return flow_chart([
        ('01','MINT','ERC-721 NFT\ncreated on-chain',BLU_H),
        ('02','REGISTER','VIN + owner\nmetadata attached',CYN_H),
        ('03','LOG','Mechanic + owner\nco-sign records',GRN_H),
        ('04','TRANSFER','Ownership\naudited on-chain',AMB_H),
        ('05','QUERY','Open API for\ninsurers & buyers',RED_H),
    ])

def ch_zk():
    return flow_chart([
        ('01','COLLECT','GPS · accel\ngyro · ECU data',BLU_H),
        ('02','COMPUTE','Off-chain\nbehavior scoring',CYN_H),
        ('03','PROVE','ZK proof\ngenerated off-chain',AMB_H),
        ('04','SUBMIT','Proof verified\nby smart contract',GRN_H),
        ('05','UPDATE','Safety score &\npremium updated',GRN_H),
    ])

def ch_claim():
    return flow_chart([
        ('1','INCIDENT','Accident,\ntheft or damage',RED_H),
        ('2','SUBMIT','Rider submits\nvia app',BLU_H),
        ('3','ORACLE\nVERIFY','Multi-source\nevidence aggregated',CYN_H),
        ('4','CONSENSUS','Validators\nreach quorum',AMB_H),
        ('5','PAYOUT','USDC released\nfrom pool',GRN_H),
        ('6','REBALANCE','Reserves &\npremiums reset',GRN_H),
    ], bw=0.13, gap=0.025)

def ch_market_bars():
    fig,ax=plt.subplots(figsize=(7,3.5),facecolor=BG_H)
    ax.set_facecolor(BG2_H)
    cats=['Global Motorcycle\nMarket','Motorcycle\nInsurance Market',
          'Value Lost to\nFraud / Mispricing','Recoverable\nValue (1-2%)']
    vals=[120,70,6,1.5]
    cols=[BLU_H,GRN_H,RED_H,AMB_H]
    bars=ax.barh(cats,vals,color=cols,height=0.5,alpha=0.9)
    for bar,v in zip(bars,vals):
        ax.text(v+0.5,bar.get_y()+bar.get_height()/2,
                f'${v}B+',va='center',fontsize=10,fontweight='bold',color=TXT_H)
    ax.set_xlim(0,145); ax.set_xlabel('$ Billions',fontsize=9,color=MUT_H)
    ax.tick_params(axis='y',labelsize=9,labelcolor=TXT_H)
    ax.tick_params(axis='x',labelsize=8,labelcolor=MUT_H)
    ax.spines[['top','right']].set_visible(False)
    ax.spines[['left','bottom']].set_color(LIN_H)
    ax.grid(axis='x',alpha=0.2,color=MUT_H,ls='--')
    ax.set_title('Market Scale of the Inefficiency',color=TXT_H,fontsize=10,pad=6)
    fig.tight_layout(); return f2s(fig)

def ch_revenue_pie():
    fig,ax=plt.subplots(figsize=(5.2,4.4),subplot_kw=dict(aspect='equal'),facecolor=BG_H)
    ax.set_facecolor(BG_H)
    labels=['Insurance Pool\nFee (44%)','API / Data\nAccess (22%)',
            'VIN Registration\n(22%)','Ownership\nTransfer (7%)',
            'Service Logging\n(4%)','Claims Processing\n(1%)']
    sizes=[44,22,22,7,4,1]; cols=[BLU_H,GRN_H,CYN_H,AMB_H,RED_H,MUT_H]
    wedges,_=ax.pie(sizes,colors=cols,startangle=90,
                    wedgeprops=dict(lw=2,edgecolor=BG_H),explode=[0.03]*6)
    ax.legend(wedges,labels,loc='lower center',bbox_to_anchor=(0.5,-0.22),
              ncol=2,fontsize=7.5,frameon=False,labelcolor=TXT_H)
    ax.set_title('Year 1 Revenue Mix',color=TXT_H,fontsize=11,fontweight='bold',pad=8)
    fig.tight_layout(); return f2s(fig)

def ch_token_donut():
    fig,ax=plt.subplots(figsize=(5.2,4.6),subplot_kw=dict(aspect='equal'),facecolor=BG_H)
    ax.set_facecolor(BG_H)
    labels=['Community &\nEcosystem 30%','Team &\nFounders 18%',
            'Treasury /\nDAO 15%','Pool\nSubsidy 15%',
            'Investors\nPrivate 9%','Investors\nSeed 8%','Liquidity 5%']
    sizes=[30,18,15,15,9,8,5]; cols=[BLU_H,GRN_H,CYN_H,AMB_H,RED_H,MUT_H,PUR_H]
    wedges,_=ax.pie(sizes,colors=cols,startangle=90,
                    wedgeprops=dict(width=0.55,lw=2,edgecolor=BG_H))
    ax.text(0,0,'1 B\n$MOTO',ha='center',va='center',
            fontsize=13,fontweight='bold',color=TXT_H,linespacing=1.6)
    ax.legend(wedges,labels,loc='lower center',bbox_to_anchor=(0.5,-0.20),
              ncol=2,fontsize=7.5,frameon=False,labelcolor=TXT_H)
    ax.set_title('$MOTO Allocation · Fixed Supply',color=TXT_H,fontsize=11,fontweight='bold',pad=8)
    fig.tight_layout(); return f2s(fig)

def ch_premium_pie():
    fig,ax=plt.subplots(figsize=(5.2,4.4),subplot_kw=dict(aspect='equal'),facecolor=BG_H)
    ax.set_facecolor(BG_H)
    labels=['Liquid Reserve\n25%','Catastrophe\nReserve 20%',
            'Safe-Rider\nRebates 18%','Expected\nClaims 12%',
            'Reinsurance\n10%','Treasury\n10%','Protocol\nFee 5%']
    sizes=[25,20,18,12,10,10,5]; cols=[GRN_H,BLU_H,CYN_H,RED_H,AMB_H,MUT_H,PUR_H]
    wedges,_=ax.pie(sizes,colors=cols,startangle=90,
                    wedgeprops=dict(lw=2,edgecolor=BG_H),explode=[0.03]*7)
    ax.legend(wedges,labels,loc='lower center',bbox_to_anchor=(0.5,-0.22),
              ncol=2,fontsize=7.5,frameon=False,labelcolor=TXT_H)
    ax.set_title('Where Every Premium Goes',color=TXT_H,fontsize=11,fontweight='bold',pad=8)
    fig.tight_layout(); return f2s(fig)

def ch_mechanic_defense():
    fig,ax=plt.subplots(figsize=(12.6,3.8),facecolor=BG_H)
    ax.set_facecolor(BG_H); ax.axis('off')
    layers=[
        ('PREVENTION',GRN_H,[
            ('Multi-party co-signing',
             'Every record needs cryptographic signatures\nfrom BOTH mechanic AND bike owner.\nEliminates all unilateral mechanic fraud.'),
            ('Hardware-attested telemetry',
             'OBD-II / Bluetooth modules read mileage\ndirectly from the bike. No human discretion\nfor the most fraud-prone data categories.'),
        ]),
        ('DETECTION',AMB_H,[
            ('Bounty market',
             'Any third party — auditors, dealers, bots —\ncan submit fraud evidence and earn a share\nof the slashed stake. Detection is incentivized.'),
            ('Trigger-on-monetization',
             'Re-verification fires automatically at every\ndownstream event: resale, financing, insurance.\nFraud is caught exactly when it would pay off.'),
        ]),
        ('ENFORCEMENT',RED_H,[
            ('Tiered staking',
             'Write access is earned through clean\noperating history. More volume = more stake\nat risk. Cannot buy credibility.'),
            ('Risk-tiered escrow + per-entry stake',
             'Fees held 30–90 days. Disputes claw back\nfees and slash stake. 50 open entries = 50×\nper-entry stake at risk — not just base stake.'),
        ]),
    ]
    cw=0.295; gap=0.035; total=3*cw+2*gap; sx=(1-total)/2
    for i,(title,col,items) in enumerate(layers):
        x=sx+i*(cw+gap)
        ax.add_patch(FancyBboxPatch((x,0.84),cw,0.14,boxstyle="round,pad=0.01",
            fc=col+'33',ec=col,lw=2,transform=ax.transAxes,clip_on=False))
        ax.text(x+cw/2,0.91,title,transform=ax.transAxes,
                ha='center',va='center',fontsize=11,fontweight='bold',color=col)
        ax.add_patch(FancyBboxPatch((x,0.02),cw,0.80,boxstyle="round,pad=0.01",
            fc=col+'0c',ec=col+'44',lw=1,transform=ax.transAxes,clip_on=False))
        y=0.77
        for item_title,item_body in items:
            ax.text(x+0.015,y,item_title,transform=ax.transAxes,
                    ha='left',va='top',fontsize=9.5,fontweight='bold',color=col)
            y-=0.13
            ax.text(x+0.015,y,item_body,transform=ax.transAxes,
                    ha='left',va='top',fontsize=7.8,color=TXT_H,linespacing=1.45)
            y-=0.30
    ax.text(0.5,0.0,'Each layer compensates for the others\' weaknesses — together they make fraud economically irrational',
            transform=ax.transAxes,ha='center',va='bottom',
            fontsize=8,color=MUT_H,style='italic')
    return f2s(fig)

def ch_pnl():
    fig,ax=plt.subplots(figsize=(9,4),facecolor=BG_H)
    ax.set_facecolor(BG2_H)
    years=['Year 1','Year 2','Year 3']
    rev=[275,1515,6448]; opex=[-1800,-2500,-3400]; net=[-1525,-985,3048]
    x=np.arange(3); w=0.24
    ax.bar(x-w,rev,w,label='Revenue (₩M)',color=GRN_H,alpha=0.9,zorder=3)
    ax.bar(x,opex,w,  label='OpEx (₩M)',   color=RED_H,alpha=0.9,zorder=3)
    ax.bar(x+w,net,w, label='Net (₩M)',    color=BLU_H,alpha=0.9,zorder=3)
    for i,(r,e,n) in enumerate(zip(rev,opex,net)):
        ax.text(i-w,r+70,f'₩{r:,}', ha='center',fontsize=7.5,color=GRN_H,fontweight='bold')
        ax.text(i,  e-70,f'₩{e:,}', ha='center',fontsize=7.5,color=RED_H,fontweight='bold',va='top')
        va='bottom' if n>=0 else 'top'
        ax.text(i+w,n+(70 if n>=0 else -70),f'₩{n:,}',ha='center',
                fontsize=7.5,color=BLU_H,fontweight='bold',va=va)
    ax.axhline(0,color=MUT_H,lw=0.8,ls='--',zorder=2)
    ax.annotate('Break-even\nYear 3',xy=(2+w,3048),xytext=(2.25,2100),
                arrowprops=dict(arrowstyle='->',color=GRN_H,lw=1.5),
                color=GRN_H,fontsize=9,fontweight='bold')
    ax.set_xticks(x); ax.set_xticklabels(years,fontsize=11,color=TXT_H)
    ax.set_ylabel('₩ Millions',color=MUT_H,fontsize=9)
    ax.tick_params(axis='y',labelcolor=MUT_H,labelsize=8)
    ax.spines[['top','right']].set_visible(False)
    ax.spines[['left','bottom']].set_color(LIN_H)
    ax.grid(axis='y',alpha=0.2,color=MUT_H,ls='--',zorder=1)
    ax.legend(fontsize=9,frameon=False,labelcolor=TXT_H,loc='upper left')
    ax.set_title('3-Year Stressed P&L · 15% claim freq · ₩2.5M avg severity',
                 color=TXT_H,fontsize=10,pad=8)
    fig.tight_layout(); return f2s(fig)

def ch_lp():
    fig,ax=plt.subplots(figsize=(7.5,3.8),facecolor=BG_H)
    ax.set_facecolor(BG2_H)
    years=['Year 1','Year 2','Year 3']
    fnd=[1500,1500,1500]; defi=[750,4000,15000]
    inst=[500,3000,12000]; strat=[250,2000,8000]
    x=np.arange(3); w=0.45
    ax.bar(x,fnd, w,label='Foundation',color=BLU_H,alpha=0.9)
    ax.bar(x,defi,w,bottom=fnd,label='DeFi LPs',color=GRN_H,alpha=0.9)
    bot3=[fnd[i]+defi[i] for i in range(3)]
    ax.bar(x,inst, w,bottom=bot3,label='Institutional',color=AMB_H,alpha=0.9)
    bot4=[bot3[i]+inst[i] for i in range(3)]
    ax.bar(x,strat,w,bottom=bot4,label='Strategic',color=CYN_H,alpha=0.9)
    for i,tot in enumerate([3000,10500,36500]):
        top=bot4[i]+strat[i]
        ax.text(i,top+400,f'₩{tot:,}M',ha='center',fontsize=9,color=TXT_H,fontweight='bold',linespacing=1.4)
    ax.set_xticks(x); ax.set_xticklabels(years,fontsize=11,color=TXT_H)
    ax.set_ylabel('₩ Millions',color=MUT_H,fontsize=9)
    ax.tick_params(axis='y',labelcolor=MUT_H,labelsize=8)
    ax.spines[['top','right']].set_visible(False)
    ax.spines[['left','bottom']].set_color(LIN_H)
    ax.grid(axis='y',alpha=0.2,color=MUT_H,ls='--')
    ax.legend(fontsize=8.5,frameon=False,labelcolor=TXT_H,loc='upper left')
    ax.set_title('LP Capital Plan by Source (₩ Millions)',color=TXT_H,fontsize=10,pad=8)
    fig.tight_layout(); return f2s(fig)

def ch_competitive():
    fig,ax=plt.subplots(figsize=(12,3.6),facecolor=BG_H)
    ax.set_facecolor(BG2_H)
    comps=['Centralized History\n(CARFAX, Korea Auto)',
           'Traditional Insurers\n(Samsung F&M, Hyundai M&F)',
           'DeFi Insurance\n(Nexus Mutual, InsurAce)',
           'MotoChain & RideTrue  ←']
    crits=['Portable\nIdentity','Behavior-\nBased Price','On-Chain\nClaims',
           'Economic\nEnforcement','Consumer\nUX','Real-World\nData Moat']
    sc=[[0,0,0,0,1,2],[0,1,0,0,2,1],[1,0,2,1,0,0],[2,2,2,2,2,2]]
    sym={0:'✗',1:'∼',2:'✓'}; fc={0:RED_H+'44',1:AMB_H+'44',2:GRN_H+'44'}; tc={0:RED_H,1:AMB_H,2:GRN_H}
    for ri,row in enumerate(sc):
        for ci,s in enumerate(row):
            ax.add_patch(mpatches.FancyBboxPatch((ci+0.06,ri+0.06),0.88,0.88,
                boxstyle="round,pad=0.05",fc=fc[s],ec='none'))
            ax.text(ci+0.5,ri+0.5,sym[s],ha='center',va='center',
                    fontsize=17,fontweight='bold',color=tc[s])
    ax.add_patch(mpatches.FancyBboxPatch((-0.06,3.02),len(crits)+0.12,0.96,
        boxstyle="round,pad=0.02",fc='none',ec=BLU_H,lw=2.5))
    ax.set_xlim(0,len(crits)); ax.set_ylim(0,len(comps))
    ax.set_xticks([i+0.5 for i in range(len(crits))])
    ax.set_xticklabels(crits,fontsize=8.5,color=TXT_H,linespacing=1.35)
    ax.set_yticks([i+0.5 for i in range(len(comps))])
    ax.set_yticklabels(comps,fontsize=9,color=TXT_H)
    ax.tick_params(length=0); ax.spines[:].set_visible(False)
    ax.set_title('Competitive Feature Matrix',color=TXT_H,fontsize=10.5,pad=8)
    els=[Line2D([0],[0],marker='s',color='w',mfc=c,ms=10,label=l)
         for c,l in [(GRN_H,'✓ Yes'),(AMB_H,'∼ Partial'),(RED_H,'✗ No')]]
    ax.legend(handles=els,loc='upper right',fontsize=8.5,frameon=False,
              labelcolor=TXT_H,bbox_to_anchor=(1,-0.1),ncol=3)
    fig.tight_layout(); return f2s(fig)

def ch_security():
    fig,ax=plt.subplots(figsize=(12.6,3.2),facecolor=BG_H)
    ax.set_facecolor(BG_H); ax.axis('off')
    phases=[
        ('Pre-Launch\nEngineering Gate',BLU_H,[
            '≥90% unit test coverage','Invariant & fuzz testing',
            'Zero unresolved critical findings','Internal threat model review']),
        ('External Audits\n(Min. 2 Firms)',CYN_H,[
            'OpenZeppelin / Trail of Bits / Spearbit',
            'Formal verification of fund flows',
            'All reports published in full',
            'Unresolved findings disclosed']),
        ('Continuous Bug\nBounty (Immunefi)',AMB_H,[
            'Critical: up to $250,000',
            'High: up to $75,000',
            'Medium: up to $15,000',
            'Live for lifetime of protocol']),
        ('Operational\nSecurity',GRN_H,[
            '5-of-9 multi-sig treasury',
            '48-hour time-locked upgrades',
            'Emergency pause (narrow scope)',
            'Annual pen tests + incident runbook']),
    ]
    n=len(phases); cw=0.215; gap=0.028; total=n*cw+(n-1)*gap; sx=(1-total)/2
    ax.axhline(0.72,xmin=sx,xmax=sx+total,color=MUT_H,lw=1.5)
    for i,(title,col,pts) in enumerate(phases):
        x=sx+i*(cw+gap); cx=x+cw/2
        ax.plot(cx,0.72,'o',color=col,ms=14,transform=ax.transAxes,zorder=4)
        ax.text(cx,0.86,title,transform=ax.transAxes,
                ha='center',va='bottom',fontsize=8.5,fontweight='bold',
                color=col,linespacing=1.4)
        ax.add_patch(FancyBboxPatch((x,0.02),cw,0.58,boxstyle="round,pad=0.01",
            fc=col+'18',ec=col+'55',lw=1,transform=ax.transAxes,clip_on=False))
        for j,pt in enumerate(pts):
            ax.text(x+0.01,0.56-j*0.125,f'• {pt}',transform=ax.transAxes,
                    ha='left',va='top',fontsize=7.5,color=TXT_H)
    return f2s(fig)

def ch_oracle():
    """Multi-source oracle aggregation diagram."""
    fig,ax=plt.subplots(figsize=(8,3.8),facecolor=BG_H)
    ax.set_facecolor(BG_H); ax.axis('off')
    ax.set_xlim(0,10); ax.set_ylim(0,5)
    sources=[
        (1.2,4.2,'Repair Shops\n& Mechanics',BLU_H),
        (1.2,3.0,'Police /\nTheft Records',CYN_H),
        (1.2,1.8,'Insurance\nDatabases',AMB_H),
        (1.2,0.6,'OEM / OBD\nTelemetry',GRN_H),
    ]
    for sx,sy,label,col in sources:
        ax.add_patch(FancyBboxPatch((sx-0.9,sy-0.38),1.8,0.72,
            boxstyle="round,pad=0.08",fc=col+'22',ec=col,lw=1.5))
        ax.text(sx,sy,label,ha='center',va='center',fontsize=8,
                fontweight='bold',color=col,linespacing=1.35)
        ax.annotate('',xy=(3.8,2.4),xytext=(sx+0.9,sy),
                    arrowprops=dict(arrowstyle='->',color=col+'aa',lw=1.5))
    ax.add_patch(FancyBboxPatch((3.8,1.65),2.4,1.5,
        boxstyle="round,pad=0.1",fc=CYN_H+'22',ec=CYN_H,lw=2))
    ax.text(5.0,2.4,'Oracle\nConsensus\nLayer',ha='center',va='center',
            fontsize=9.5,fontweight='bold',color=CYN_H,linespacing=1.4)
    ax.annotate('',xy=(7.1,3.2),xytext=(6.2,3.0),
                arrowprops=dict(arrowstyle='->',color=GRN_H,lw=1.8))
    ax.add_patch(FancyBboxPatch((7.1,2.6),2.1,1.2,
        boxstyle="round,pad=0.08",fc=GRN_H+'22',ec=GRN_H,lw=2))
    ax.text(8.15,3.2,'Smart\nContract\nPayout',ha='center',va='center',
            fontsize=9,fontweight='bold',color=GRN_H,linespacing=1.4)
    ax.annotate('',xy=(7.1,1.6),xytext=(6.2,2.0),
                arrowprops=dict(arrowstyle='->',color=RED_H,lw=1.8,
                                linestyle='dashed'))
    ax.add_patch(FancyBboxPatch((7.1,0.9),2.1,1.2,
        boxstyle="round,pad=0.08",fc=RED_H+'22',ec=RED_H,lw=2))
    ax.text(8.15,1.5,'Arbitration\n(disputed\nclaims)',ha='center',va='center',
            fontsize=8.5,fontweight='bold',color=RED_H,linespacing=1.4)
    ax.text(5.0,0.15,'No single source is trusted in isolation — multiple oracles must reach consensus',
            ha='center',va='center',fontsize=7.5,color=MUT_H,style='italic')
    return f2s(fig)

def ch_two_sided():
    """Two-sided market visual."""
    fig,ax=plt.subplots(figsize=(10,4),facecolor=BG_H)
    ax.set_facecolor(BG_H); ax.axis('off')
    ax.set_xlim(0,10); ax.set_ylim(0,5)
    # Protocol hub
    ax.add_patch(FancyBboxPatch((3.8,1.6),2.4,1.8,boxstyle="round,pad=0.1",
        fc=BLU_H+'22',ec=BLU_H,lw=2.5))
    ax.text(5.0,2.5,'MotoChain\n& RideTrue',ha='center',va='center',
            fontsize=10,fontweight='bold',color=BLU_H,linespacing=1.4)
    # Supply side
    supply=[('Riders',1.2,4.1,GRN_H),('DeFi LPs',1.2,2.8,GRN_H),
            ('Mechanics',1.2,1.5,GRN_H),('Dealers',1.2,0.5,GRN_H)]
    for label,x,y,col in supply:
        ax.add_patch(FancyBboxPatch((x-0.85,y-0.28),1.7,0.56,
            boxstyle="round,pad=0.07",fc=col+'1a',ec=col,lw=1.5))
        ax.text(x,y,label,ha='center',va='center',fontsize=9,fontweight='bold',color=col)
        ax.annotate('',xy=(3.8,2.5),xytext=(x+0.85,y),
                    arrowprops=dict(arrowstyle='->',color=col+'99',lw=1.4))
    # Demand side
    demand=[('Insurers',8.8,4.1,AMB_H),('Lenders',8.8,2.8,AMB_H),
            ('Buyers',8.8,1.5,AMB_H),('Marketplaces',8.8,0.5,AMB_H)]
    for label,x,y,col in demand:
        ax.add_patch(FancyBboxPatch((x-0.85,y-0.28),1.7,0.56,
            boxstyle="round,pad=0.07",fc=col+'1a',ec=col,lw=1.5))
        ax.text(x,y,label,ha='center',va='center',fontsize=9,fontweight='bold',color=col)
        ax.annotate('',xy=(x-0.85,y),xytext=(6.2,2.5),
                    arrowprops=dict(arrowstyle='->',color=col+'99',lw=1.4))
    ax.text(5.0,4.7,'SUPPLY SIDE',ha='center',fontsize=8.5,fontweight='bold',color=GRN_H,transform=ax.transAxes)
    ax.text(1.0,4.7,'◀  Capital & Data',ha='center',fontsize=8,color=GRN_H,style='italic',transform=ax.transAxes)
    ax.text(5.0,4.7,'DEMAND SIDE',ha='center',fontsize=8.5,fontweight='bold',color=AMB_H)
    ax.text(0.5,4.8,'SUPPLY SIDE',ha='center',fontsize=9,fontweight='bold',color=GRN_H)
    ax.text(9.5,4.8,'DEMAND SIDE',ha='center',fontsize=9,fontweight='bold',color=AMB_H)
    return f2s(fig)

# ── PPTX helpers ───────────────────────────────────────────────────────────
def new_prs():
    prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH; return prs

def blank(prs):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb=BG_R; return s

def img(slide,stream,l,t,w,h=None):
    stream.seek(0)
    if h: slide.shapes.add_picture(stream,l,t,w,h)
    else: slide.shapes.add_picture(stream,l,t,w)

def tb(slide,text,l,t,w,h,sz=11,bold=False,col=TXT_R,
       align=PP_ALIGN.LEFT,italic=False,wrap=True):
    box=slide.shapes.add_textbox(l,t,w,h)
    tf=box.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold
    r.font.italic=italic; r.font.color.rgb=col
    return box

def multiline_tb(slide,lines,l,t,w,h):
    """lines = [(text, size, bold, color, space_before)]"""
    box=slide.shapes.add_textbox(l,t,w,h)
    tf=box.text_frame; tf.word_wrap=True
    for i,(text,sz,bold,col,sp) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_before=Pt(sp)
        r=p.add_run(); r.text=text
        r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=col
    return box

def card(slide,l,t,w,h,fill=CARD_H,border=None,border_w=1.5):
    s=slide.shapes.add_shape(1,l,t,w,h)
    s.fill.solid(); s.fill.fore_color.rgb=_c(fill)
    if border: s.line.color.rgb=_c(border); s.line.width=Pt(border_w)
    else: s.line.fill.background()
    return s

def pagenum(slide,n,total):
    tb(slide,f'{n} / {total}',
       SW-Inches(1.08),SH-Inches(0.34),Inches(1.0),Inches(0.28),
       sz=8,col=MUT_R,align=PP_ALIGN.RIGHT)

def slide_header(slide, section_label):
    """Branded top bar on every non-title slide."""
    # Bar background
    bar=slide.shapes.add_shape(1,Inches(0),Inches(0),SW,HDR_H)
    bar.fill.solid(); bar.fill.fore_color.rgb=BG2_R
    bar.line.color.rgb=LIN_R; bar.line.width=Pt(1)
    # Brand left
    tb(slide,'MotoChain & RideTrue',
       Inches(0.35),Inches(0.06),Inches(5),Inches(0.3),
       sz=9.5,bold=True,col=BLU_R)
    # Authors right
    tb(slide,'Mireu Kim  ·  Enxhi Topalli',
       SW-Inches(3.2),Inches(0.06),Inches(3.0),Inches(0.3),
       sz=8.5,col=MUT_R,align=PP_ALIGN.RIGHT)
    # Section label
    tb(slide,section_label,
       Inches(0.35),Inches(0.23),SW-Inches(0.7),Inches(0.22),
       sz=7.5,col=_c(AMB_H))

def eyebrow(slide,text):
    tb(slide,text,ML,EYEBROW_Y,CW,Inches(0.26),sz=8,bold=True,col=BLU_R)

def heading(slide,text,sz=26,col=TXT_R):
    tb(slide,text,ML,TITLE_Y,CW,Inches(0.62),sz=sz,bold=True,col=col)

def subheading(slide,text,col=MUT_R):
    tb(slide,text,ML,TITLE_Y+Inches(0.6),CW,Inches(0.32),sz=10.5,italic=True,col=col)

def bullet_box(slide,title,points,l,t,w,h,title_col=TXT_R,pt_col=TXT_R,title_sz=11,pt_sz=10):
    box=slide.shapes.add_textbox(l,t,w,h)
    tf=box.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; r=p.add_run()
    r.text=title; r.font.size=Pt(title_sz); r.font.bold=True; r.font.color.rgb=title_col
    for pt in points:
        p2=tf.add_paragraph(); p2.space_before=Pt(5)
        r2=p2.add_run(); r2.text=f'• {pt}'
        r2.font.size=Pt(pt_sz); r2.font.color.rgb=pt_col

# ── SLIDE BUILDERS ─────────────────────────────────────────────────────────

def s01_title(slide):
    # Large title
    tb(slide,'MotoChain & RideTrue',
       Inches(0.7),Inches(1.5),Inches(9),Inches(1.3),
       sz=44,bold=True,col=TXT_R)
    tb(slide,'A Verifiable Identity & Behavior-Based\nInsurance Protocol for Motorcycles',
       Inches(0.7),Inches(2.95),Inches(9),Inches(0.95),
       sz=17,col=MUT_R)
    # Accent line
    bar=slide.shapes.add_shape(1,Inches(0.7),Inches(3.95),Inches(7),Pt(2.5))
    bar.fill.solid(); bar.fill.fore_color.rgb=BLU_R; bar.line.fill.background()
    tb(slide,'Mireu Kim  ·  Enxhi Topalli',
       Inches(0.7),Inches(4.18),Inches(8),Inches(0.38),sz=13,col=BLU_R)
    tb(slide,'CS495 DeFi  ·  May 2025',
       Inches(0.7),Inches(4.62),Inches(8),Inches(0.32),sz=11,col=MUT_R)
    # Right KPI boxes
    for i,(kpi,label,col) in enumerate([
        ('$120 B+','Global Motorcycle Market',BLU_H),
        ('$70 B+', 'Motorcycle Insurance Market',GRN_H),
        ('Base L2', 'Ethereum Layer 2 Network',CYN_H)]):
        x=Inches(9.9); y=Inches(2.0+i*1.6)
        card(slide,x,y,Inches(3.1),Inches(1.35),CARD_H,col,1.5)
        tb(slide,kpi,  x+Inches(0.18),y+Inches(0.1), Inches(2.7),Inches(0.6), sz=24,bold=True,col=_c(col))
        tb(slide,label,x+Inches(0.18),y+Inches(0.75),Inches(2.7),Inches(0.42),sz=9.5,col=MUT_R)

def s02_pitch(slide):
    slide_header(slide,"I. THE ELEVATOR PITCH")
    eyebrow(slide,'I. THE ELEVATOR PITCH')
    heading(slide,'A Dual-Protocol DeFi System for Trust')
    subheading(slide,'On-chain identity + economic incentives + ZK proofs = verifiable, self-enforcing markets')
    t=BODY_Y+Inches(0.1)
    bw=(CW-Inches(0.25))/2
    for i,(title,col,pts) in enumerate([
        ('MotoChain',BLU_H,[
            'Tamper-proof identity layer for motorcycles',
            'VIN-linked ERC-721 NFT for every bike',
            'Immutable, staked service records',
            'Solves information asymmetry in resale',
            'Portable history owned by the user',
        ]),
        ('RideTrue',GRN_H,[
            'Behavior-based insurance market for riders',
            'ZK proofs for privacy-preserving scoring',
            'Pooled USDC insurance — claims on-chain',
            'Replaces demographic proxies with real data',
            'Safe-rider rebates from verified behavior',
        ])]):
        x=ML+i*(bw+Inches(0.25))
        card(slide,x,t,bw,Inches(4.8),CARD_H,col,1.5)
        tb(slide,title,x+Inches(0.2),t+Inches(0.15),bw,Inches(0.52),sz=22,bold=True,col=_c(col))
        bxt=slide.shapes.add_textbox(x+Inches(0.2),t+Inches(0.74),bw-Inches(0.35),Inches(3.8))
        tf=bxt.text_frame; tf.word_wrap=True
        for j,pt in enumerate(pts):
            p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
            p.space_before=Pt(7); r=p.add_run(); r.text=f'• {pt}'
            r.font.size=Pt(11); r.font.color.rgb=TXT_R
    tb(slide,'Core innovation: combining on-chain identity, economic enforcement through staking, and '
             'privacy-preserving ZK proofs — transforming trust-dependent markets into verifiable, self-enforcing systems.',
       ML,t+Inches(5.05),CW,Inches(0.42),sz=10,italic=True,col=AMB_R)

def s03_mission(slide):
    slide_header(slide,"MISSION STATEMENT")
    eyebrow(slide,'MISSION STATEMENT')
    heading(slide,'Verifiable Truth, User-Owned Identity, Economically Enforced')
    t=BODY_Y+Inches(0.1)
    # Big quote
    card(slide,ML,t,CW,Inches(1.5),CARD_H,BLU_H,1.5)
    tb(slide,'"MotoChain and RideTrue aim to eliminate trust-based inefficiencies in real-world markets\n'
             'by replacing them with verifiable, user-owned data and economically enforced truth systems."',
       ML+Inches(0.25),t+Inches(0.2),CW-Inches(0.5),Inches(1.15),
       sz=12.5,italic=True,col=TXT_R)
    # Vision stages
    stages=[
        ('Starting point','Motorcycles & rider insurance in Seoul',BLU_H),
        ('Medium term', 'Broader mobility markets — Vietnam, Indonesia, Thailand',CYN_H),
        ('Long-term vision','Global trustless layer for real-world asset verification\nand behavior-based financial services',GRN_H),
    ]
    sw2=(CW-Inches(0.4))/3
    for i,(stage,desc,col) in enumerate(stages):
        x=ML+i*(sw2+Inches(0.2))
        y=t+Inches(1.72)
        card(slide,x,y,sw2,Inches(2.7),CARD_H,col,1.2)
        tb(slide,stage,x+Inches(0.15),y+Inches(0.15),sw2,Inches(0.35),sz=10,bold=True,col=_c(col))
        tb(slide,desc, x+Inches(0.15),y+Inches(0.55),sw2-Inches(0.3),Inches(1.8),sz=10,col=TXT_R)
        if i<2:
            tb(slide,'→',x+sw2+Inches(0.02),y+Inches(1.1),Inches(0.18),Inches(0.5),sz=16,col=MUT_R,align=PP_ALIGN.CENTER)
    tb(slide,'Individuals — not institutions — control their history, reputation, and financial outcomes.',
       ML,SH-Inches(0.5),CW,Inches(0.32),sz=10.5,italic=True,col=AMB_R,align=PP_ALIGN.CENTER)

def s04_problem_market(slide,market_img):
    slide_header(slide,"II.1 — COMPREHENSIVE PROBLEM & ROOT CAUSE ANALYSIS")
    eyebrow(slide,'II.1 — PROBLEM & ROOT CAUSE ANALYSIS')
    heading(slide,'Two Broken Markets — One Structural Failure: Trust')
    t=BODY_Y
    # Market stats
    bw=Inches(2.8)
    for i,(kpi,mkt,col,desc) in enumerate([
        ('$120 B+','Global Motorcycle Market',BLU_H,
         'Sellers possess full knowledge of vehicle history.\nBuyers depend on unverifiable claims.\n→ Classic "market for lemons" (Akerlof 1970)'),
        ('$70 B+','Insurance Market',AMB_H,
         'Risk priced using age, demographics, not behavior.\nSafe riders subsidize risky ones.\n→ Systematic mispricing, reduced trust'),
    ]):
        x=ML+i*(bw+Inches(0.2))
        card(slide,x,t,bw,Inches(4.2),CARD_H,col,1.5)
        tb(slide,kpi, x+Inches(0.18),t+Inches(0.12),bw,Inches(0.78),sz=32,bold=True,col=_c(col))
        tb(slide,mkt, x+Inches(0.18),t+Inches(0.92),bw,Inches(0.35),sz=13,bold=True,col=TXT_R)
        tb(slide,desc,x+Inches(0.18),t+Inches(1.32),bw-Inches(0.3),Inches(2.5),sz=10,col=MUT_R)
    # Bar chart
    img(slide,market_img,ML+Inches(6.1),t,Inches(6.65),Inches(4.2))
    # Root cause bar
    card(slide,ML,t+Inches(4.4),CW,Inches(0.72),CARD_H,RED_H,1.5)
    tb(slide,'Root cause:  No shared, tamper-proof, portable record of real-world events — '
             'data is siloed, mutable, and inaccessible across platforms.',
       ML+Inches(0.25),t+Inches(0.52),CW-Inches(0.5),Inches(0.52),sz=12,bold=True,col=RED_R)

def s05_why_blockchain(slide):
    slide_header(slide,"II.1 — WHY BLOCKCHAIN IS A STRUCTURAL REQUIREMENT")
    eyebrow(slide,'II.1 — WHY BLOCKCHAIN IS A STRUCTURAL REQUIREMENT')
    heading(slide,'Centralized Systems Cannot Provide All Four Properties Simultaneously')
    t=BODY_Y+Inches(0.1)
    # Centralized fails
    card(slide,ML,t,Inches(5.8),Inches(4.8),CARD_H,RED_H,1.5)
    tb(slide,'Why Centralized Alternatives Fail',ML+Inches(0.2),t+Inches(0.15),Inches(5.4),Inches(0.38),
       sz=13,bold=True,col=RED_R)
    fails=[
        'Data-holders have conflicts of interest in monetizing access',
        'Competing institutions have no incentive to share data openly',
        'Centralized records remain mutable and selectively disclosable',
        'No economic enforcement — data providers have no "skin in the game"',
        'Users cannot carry their history across platforms',
    ]
    bxt=slide.shapes.add_textbox(ML+Inches(0.2),t+Inches(0.6),Inches(5.4),Inches(4.0))
    tf=bxt.text_frame; tf.word_wrap=True
    for j,f in enumerate(fails):
        p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
        p.space_before=Pt(8); r=p.add_run(); r.text=f'✗  {f}'
        r.font.size=Pt(10.5); r.font.color.rgb=TXT_R
    # 4 blockchain properties
    props=[
        ('Immutability','Once written, history cannot be altered by any single actor',BLU_H),
        ('Decentralized\nValidation','Multiple independent participants validate every entry',CYN_H),
        ('Economic\nEnforcement','Staking and slashing create real financial consequences for fraud',GRN_H),
        ('Portable\nUser-Owned Identity','Reputation and history follow the user across platforms',AMB_H),
    ]
    pw=(CW-Inches(5.8)-Inches(0.25))/2
    for i,(title,desc,col) in enumerate(props):
        r2=i//2; c2=i%2
        x=ML+Inches(6.05)+c2*(pw+Inches(0.15))
        y=t+r2*(Inches(2.35)+Inches(0.15))
        card(slide,x,y,pw,Inches(2.35),CARD_H,col,1.2)
        tb(slide,title,x+Inches(0.15),y+Inches(0.15),pw,Inches(0.55),sz=11,bold=True,col=_c(col))
        tb(slide,desc, x+Inches(0.15),y+Inches(0.75),pw-Inches(0.3),Inches(1.4),sz=10,col=TXT_R)
    tb(slide,'In this context, blockchain is not an implementation choice — it is a structural requirement.',
       ML,SH-Inches(0.45),CW,Inches(0.32),sz=10.5,italic=True,col=AMB_R,align=PP_ALIGN.CENTER)

def s06_architecture(slide,vin_img):
    slide_header(slide,"II.2 — TECHNICAL ARCHITECTURE & SYSTEM LOGIC")
    eyebrow(slide,'II.2 — TECHNICAL ARCHITECTURE & SYSTEM LOGIC')
    heading(slide,'MotoChain: VIN Token Lifecycle on Base L2')
    subheading(slide,'Each motorcycle is anchored by an ERC-721 NFT. Every event is co-signed and permanently locked on-chain.')
    img(slide,vin_img,ML,BODY_Y+Inches(0.15),CW,Inches(2.15))
    props=[
        ('Tamper-Resistant','Once written, no single actor can alter the record. Immutability is enforced at the protocol level.',BLU_H),
        ('Multi-Party Signed','Every service entry requires cryptographic signatures from BOTH the mechanic and the bike owner.',GRN_H),
        ('Auto-Validation','Smart contracts flag inconsistencies — non-sequential mileage, conflicting records — automatically.',CYN_H),
        ('Composable','Open API for insurers, lenders, and marketplaces. History is a public good, not a silo.',AMB_H),
    ]
    pw=CW/4
    for i,(title,desc,col) in enumerate(props):
        x=ML+i*pw
        card(slide,x,BODY_Y+Inches(2.5),pw-Inches(0.12),Inches(2.8),CARD_H,col,1.2)
        tb(slide,title,x+Inches(0.14),BODY_Y+Inches(2.65),pw-Inches(0.28),Inches(0.38),sz=11,bold=True,col=_c(col))
        tb(slide,desc, x+Inches(0.14),BODY_Y+Inches(3.08),pw-Inches(0.28),Inches(2.1),sz=9.5,col=TXT_R)

def s07_zk_antispoofing(slide,zk_img):
    slide_header(slide,"II.2.1 — TELEMATICS INTEGRITY & ANTI-SPOOFING")
    eyebrow(slide,'II.2.1 — TELEMATICS INTEGRITY & ANTI-SPOOFING')
    heading(slide,'RideTrue: Scoring Behavior Without Exposing Your Location')
    subheading(slide,'ZK proofs validate computation — not input authenticity. Four defenses close the remaining gap.')
    img(slide,zk_img,ML,BODY_Y+Inches(0.12),CW,Inches(2.05))
    defenses=[
        ('Hardware\nAttestation',BLU_H,
         'Android Play Integrity API + iOS DeviceCheck confirm data originates from a genuine, uncompromised device.'),
        ('Sensor Fusion',CYN_H,
         'GPS + accelerometer + gyroscope + ECU data must be internally consistent across all sensors — much harder to spoof than a single feed.'),
        ('Statistical\nAnomaly Detection',AMB_H,
         'Per-rider behavioral baseline detects physically implausible trajectories, instantaneous position changes, and accelerometer inconsistencies.'),
        ('Sampling &\nChallenge Mechanics',GRN_H,
         'Random rides re-verified against cell tower triangulation, traffic data, and (optionally) paired hardware tokens installed on the bike.'),
    ]
    dw=CW/4
    for i,(title,col,desc) in enumerate(defenses):
        x=ML+i*dw
        card(slide,x,BODY_Y+Inches(2.35),dw-Inches(0.12),Inches(2.9),CARD_H,col,1.2)
        tb(slide,title,x+Inches(0.14),BODY_Y+Inches(2.5),dw-Inches(0.28),Inches(0.52),sz=11,bold=True,col=_c(col))
        tb(slide,desc, x+Inches(0.14),BODY_Y+Inches(3.08),dw-Inches(0.28),Inches(2.1),sz=9.5,col=TXT_R)

def s08_claims_oracles(slide,claim_img,oracle_img):
    slide_header(slide,"II.2.2 — INSURANCE POOL, CLAIMS & ORACLES")
    eyebrow(slide,'II.2.2 — INSURANCE POOL, CLAIMS & ORACLES')
    heading(slide,'From Premium Deposit to Automated Payout')
    img(slide,claim_img,ML,BODY_Y,CW,Inches(1.9))
    t=BODY_Y+Inches(2.05)
    # Oracle section
    img(slide,oracle_img,ML,t,Inches(7.0),Inches(3.8))
    x2=ML+Inches(7.2); bw2=CW-Inches(7.2)
    tb(slide,'Oracle Design Principles',x2,t,bw2,Inches(0.35),sz=12,bold=True,col=CYN_R)
    principles=[
        'Multi-source aggregation — no single provider trusted in isolation',
        'Economic staking by oracle operators — bad data = slashed stake',
        'Stale or conflicting data → verification-pending state, not auto-execution',
        'High-value claims require higher consensus thresholds and longer review',
        'DeFi integration: idle capital deployed to Aave for yield on excess reserves',
    ]
    bxt=slide.shapes.add_textbox(x2,t+Inches(0.42),bw2,Inches(3.2))
    tf=bxt.text_frame; tf.word_wrap=True
    for j,p2 in enumerate(principles):
        p3=tf.paragraphs[0] if j==0 else tf.add_paragraph()
        p3.space_before=Pt(7); r=p3.add_run(); r.text=f'• {p2}'
        r.font.size=Pt(10.5); r.font.color.rgb=TXT_R

def s09_ux(slide):
    slide_header(slide,"II.3 — USER EXPERIENCE & ONBOARDING FRAMEWORK")
    eyebrow(slide,'II.3 — USER EXPERIENCE & ONBOARDING FRAMEWORK')
    heading(slide,'DeFi That Feels Like a Normal App')
    subheading(slide,'Account abstraction + gasless transactions + mobile-first — no seed phrases, no gas fees, no crypto knowledge needed.')
    t=BODY_Y+Inches(0.2)
    ifaces=[
        ('Rider Dashboard',BLU_H,[
            'Safety score, premium level, projected savings — one screen',
            'Behavioral tracking runs passively in the background',
            'Rebate eligibility and score history always visible',
            'Dynamic premium adjustments and periodic rebates',
            'Retention: financial incentives compound with safe behavior',
        ]),
        ('Mechanic Interface',GRN_H,[
            'QR code scans identify the vehicle instantly',
            'Structured form captures service type, mileage, parts',
            'Owner co-signs via rider app — QR handoff at the shop',
            'End-to-end record entry: under 30 seconds',
            'Reputation capital builds with each verified entry',
        ]),
        ('Buyer Interface',AMB_H,[
            'Scan a code or enter a VIN — no account required',
            'Complete, cryptographically verified history displayed',
            'Ownership transfers, service records, claims visible',
            'MotoChain Verified badge at point of sale',
            'No wallet, no crypto — just verifiable truth',
        ]),
    ]
    bw=(CW-Inches(0.4))/3
    for i,(title,col,pts) in enumerate(ifaces):
        x=ML+i*(bw+Inches(0.2))
        card(slide,x,t,bw,Inches(5.1),CARD_H,col,1.5)
        tb(slide,title,x+Inches(0.18),t+Inches(0.15),bw,Inches(0.42),sz=14,bold=True,col=_c(col))
        bxt=slide.shapes.add_textbox(x+Inches(0.18),t+Inches(0.65),bw-Inches(0.35),Inches(4.2))
        tf=bxt.text_frame; tf.word_wrap=True
        for j,pt in enumerate(pts):
            p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
            p.space_before=Pt(8); r=p.add_run(); r.text=f'• {pt}'
            r.font.size=Pt(10.5); r.font.color.rgb=TXT_R

def s10_audience(slide,two_sided_img):
    slide_header(slide,"II.4 — TARGET AUDIENCE & CUSTOMER ACQUISITION")
    eyebrow(slide,'II.4 — TARGET AUDIENCE & CUSTOMER ACQUISITION')
    heading(slide,'Two-Sided Marketplace — Both Sides Must Join Simultaneously')
    t=BODY_Y
    img(slide,two_sided_img,ML,t,Inches(9.0),Inches(4.0))
    x2=ML+Inches(9.2); bw2=CW-Inches(9.2)
    acq=[
        ('Seoul First',BLU_H,'Launch in dense urban motorcycle market with high-volume dealers & repair shops.'),
        ('Dealer Pull',GRN_H,'Verified history → faster sales, higher prices, fewer disputes → dealers adopt first.'),
        ('Mechanic Push',CYN_H,'Dealer adoption forces mechanic-side onboarding through commercial pressure.'),
        ('Sybil Resistance',AMB_H,'Rewards require verified actions: real vehicle, real rides, real service logs.'),
        ('B2B Scale',RED_H,'Marketplace + insurer + lender integrations convert data into immediate value.'),
    ]
    for i,(label,col,desc) in enumerate(acq):
        y=t+Inches(i*0.78)
        card(slide,x2,y,bw2,Inches(0.68),CARD_H,col,1.0)
        tb(slide,label,x2+Inches(0.12),y+Inches(0.06),bw2,Inches(0.26),sz=9,bold=True,col=_c(col))
        tb(slide,desc, x2+Inches(0.12),y+Inches(0.33),bw2-Inches(0.15),Inches(0.3),sz=8.5,col=TXT_R)
    tb(slide,'Expansion: Seoul → Vietnam, Indonesia, Thailand — Asia\'s largest two-wheeler markets.',
       ML,t+Inches(4.2),CW,Inches(0.32),sz=10,italic=True,col=AMB_R)

def s11_revenue(slide,rev_img):
    slide_header(slide,"II.5 — ECONOMIC ENGINEERING & TOKENOMICS  ·  II.5.1–5.2 REVENUE STREAMS")
    eyebrow(slide,'II.5.1-5.2 — ECONOMIC DESIGN & REVENUE STREAMS')
    heading(slide,'Revenue Built on Real Usage — Not Speculative Token Emissions')
    subheading(slide,'The token is not the product. The token is an enforcement and coordination tool.')
    t=BODY_Y+Inches(0.2)
    img(slide,rev_img,Inches(0.1),t-Inches(0.1),Inches(5.8),Inches(5.6))
    x=ML+Inches(5.6); bw=CW-Inches(5.6)
    streams=[
        ('VIN Registration','₩10,000 / bike  ·  5,000 bikes Y1  →  ₩50 M',BLU_H),
        ('Ownership Transfer','₩15,000 / transfer  ·  1,000 Y1  →  ₩15 M',CYN_H),
        ('Service Logging','₩5,000 / record  ·  10,000 Y1  →  ₩50 M',GRN_H),
        ('Insurance Pool Fee','5% of premium pool  ·  ₩2B pool  →  ₩100 M',AMB_H),
        ('Claims Processing','1.5% of claim value  ·  1% net margin  →  ₩2.25 M',RED_H),
        ('API / Data Access','Subscription  ·  Highest long-term margin  →  ₩50 M',MUT_H),
    ]
    for i,(name,detail,col) in enumerate(streams):
        ry=t+Inches(i*0.78)
        card(slide,x,ry,bw,Inches(0.68),CARD_H,col,1.0)
        tb(slide,name,  x+Inches(0.14),ry+Inches(0.07),bw,Inches(0.28),sz=10,bold=True,col=_c(col))
        tb(slide,detail,x+Inches(0.14),ry+Inches(0.37),bw-Inches(0.2),Inches(0.26),sz=9,col=TXT_R)
    total_y=t+Inches(len(streams)*0.78)
    card(slide,x,total_y,bw,Inches(0.55),CARD_H,GRN_H,2.0)
    tb(slide,'TOTAL YEAR 1  →  ₩ 227 M',x+Inches(0.14),total_y+Inches(0.1),bw,Inches(0.35),sz=13,bold=True,col=GRN_R)

def s12_financials(slide):
    slide_header(slide,"II.5.3–5.4 — CORE FINANCIAL ASSUMPTIONS · SEOUL LAUNCH")
    eyebrow(slide,'II.5.3-5.4 — FINANCIAL ASSUMPTIONS & REVENUE DETAIL')
    heading(slide,'Year 1 Revenue Model: Seoul Launch · Conservative Loss Ratio (~8.9%)')
    t=BODY_Y
    # Assumptions table left
    card(slide,ML,t,Inches(6.1),Inches(5.5),CARD_H,BLU_H,1.2)
    tb(slide,'Core Assumptions (Seoul Y1)',ML+Inches(0.2),t+Inches(0.15),Inches(5.7),Inches(0.35),
       sz=12,bold=True,col=BLU_R)
    rows1=[
        ('Registered motorcycles','5,000'),('Monthly insurance riders','2,000'),
        ('Avg annual premium','₩1,080,000'),('Monthly premium pool','₩180M'),
        ('Annual premium pool','₩2,160M'),('Modeled loss ratio','~8.9% (intentionally conservative)'),
        ('Avg monthly premium','₩90,000'),('Claim frequency','8%'),
        ('Avg claim severity','₩1,200,000'),('Protocol fee','5% of premiums'),
    ]
    for j,(label,val) in enumerate(rows1):
        ry=t+Inches(0.65+j*0.45)
        bg=CARD_H if j%2==0 else LIN_H
        card(slide,ML+Inches(0.12),ry,Inches(5.86),Inches(0.4),bg)
        tb(slide,label,ML+Inches(0.25),ry+Inches(0.07),Inches(3.5),Inches(0.28),sz=9.5,col=TXT_R)
        tb(slide,val,  ML+Inches(3.8),ry+Inches(0.07),Inches(2.2),Inches(0.28),sz=9.5,bold=True,col=BLU_R)
    # Revenue detail right
    x2=ML+Inches(6.35); bw2=CW-Inches(6.35)
    card(slide,x2,t,bw2,Inches(5.5),CARD_H,GRN_H,1.2)
    tb(slide,'Revenue Stream Detail',x2+Inches(0.2),t+Inches(0.15),bw2,Inches(0.35),
       sz=12,bold=True,col=GRN_R)
    rows2=[
        ('VIN Registration','5,000 × ₩10,000','₩50 M'),
        ('Ownership Transfer','1,000 × ₩15,000','₩15 M'),
        ('Service Logging (protocol share)','10,000 × ₩5,000 × 20%','₩10 M'),
        ('Insurance Pool Fee','₩2,160M × 5%','₩108 M'),
        ('Claims Processing Fee','₩150M × 1.5%','₩2.25 M'),
        ('API / Data Access','Pilot subscriptions','₩50 M'),
        ('TOTAL','','₩235 M'),
    ]
    for j,(stream,calc,rev) in enumerate(rows2):
        ry=t+Inches(0.65+j*0.67)
        bold_row=(j==len(rows2)-1)
        bg=CARD_H if j%2==0 else LIN_H
        if bold_row: bg='#0d2010'
        card(slide,x2+Inches(0.12),ry,bw2-Inches(0.24),Inches(0.62),bg,
             border=GRN_H if bold_row else None,border_w=1.5 if bold_row else 0)
        col=GRN_R if bold_row else TXT_R
        tb(slide,stream,x2+Inches(0.25),ry+Inches(0.08),Inches(2.5),Inches(0.25),sz=9.5,bold=bold_row,col=col)
        tb(slide,calc,  x2+Inches(0.25),ry+Inches(0.33),Inches(2.5),Inches(0.22),sz=8,col=MUT_R)
        tb(slide,rev,   x2+bw2-Inches(1.2),ry+Inches(0.12),Inches(1.0),Inches(0.35),sz=11 if bold_row else 10,bold=True,col=GRN_R,align=PP_ALIGN.RIGHT)
    tb(slide,'Note: Loss ratio of ~8.9% is intentionally conservative. Industry-standard motorcycle insurance loss ratios are 55–70%. '
             'Protocol prioritizes pool solvency over premium competitiveness during first 12–24 months.',
       ML,SH-Inches(0.5),CW,Inches(0.38),sz=8.5,italic=True,col=AMB_R)

def s13_token(slide,tok_img):
    slide_header(slide,"II.5.5 — TOKEN DESIGN ($MOTO) — DRAFT ASSUMPTIONS")
    eyebrow(slide,'II.5.5 — $MOTO TOKEN DESIGN · DRAFT ASSUMPTIONS')
    heading(slide,'$MOTO: Utility & Governance Only — Not a Speculative Yield Engine')
    t=BODY_Y
    img(slide,tok_img,Inches(0.1),t-Inches(0.1),Inches(5.8),Inches(5.8))
    x=ML+Inches(5.6); bw=CW-Inches(5.6)
    tb(slide,'Token Utility',x,t,bw,Inches(0.32),sz=12,bold=True,col=TXT_R)
    utils=[
        ('Staking & Slashing',GRN_H,
         'Mechanics, validators, oracle operators stake $MOTO for write access. Misbehavior is penalized via slashing.'),
        ('Governance',BLU_H,
         'Token holders vote on fees, reserve ratios, rebate formulas, treasury deployment, and integrations.'),
        ('Reputation Gating',AMB_H,
         'Tiered staking unlocks higher write limits and reward shares for established mechanics and dealers.'),
        ('Fee Discounts',CYN_H,
         'Select API and registration fees payable in $MOTO at a discount. Received tokens burned or accrued by treasury.'),
    ]
    for i,(title,col,desc) in enumerate(utils):
        uy=t+Inches(0.38+i*1.2)
        card(slide,x,uy,bw,Inches(1.08),CARD_H,col,1.2)
        tb(slide,title,x+Inches(0.14),uy+Inches(0.1),bw,Inches(0.32),sz=10.5,bold=True,col=_c(col))
        tb(slide,desc, x+Inches(0.14),uy+Inches(0.46),bw-Inches(0.28),Inches(0.55),sz=9.5,col=TXT_R)
    tb(slide,'⚠  Critical: Insurance pool returns, rider rebates, and all protocol revenue are paid in USDC — NOT in $MOTO. '
             'This is intentional: preserves economic discipline and reduces risk of securities classification.',
       x,t+Inches(5.3),bw,Inches(0.58),sz=9,italic=True,col=AMB_R)
    tb(slide,'Network: Base L2  ·  Supply: 1,000,000,000 fixed  ·  Standard: ERC-20 with governance + staking modules',
       x,t+Inches(5.95),bw,Inches(0.28),sz=8.5,col=MUT_R)

def s14_mechanic(slide,mech_img):
    slide_header(slide,"II.5.6 — MECHANIC TRUST: A LAYERED DEFENSE")
    eyebrow(slide,'II.5.6 — MECHANIC TRUST: A LAYERED DEFENSE')
    heading(slide,'Prevention → Detection → Enforcement')
    subheading(slide,'Without trustworthy service records, the vehicle history layer is no better than a centralized database.')
    img(slide,mech_img,ML,BODY_Y+Inches(0.15),CW,Inches(4.3))
    # Tier table
    t2=BODY_Y+Inches(4.65)
    cols_h=['Tier','Type','Stake Required','Writes / Month','Fee Share','Unlock Requirement']
    widths=[Inches(0.4),Inches(1.6),Inches(1.3),Inches(1.4),Inches(0.9),Inches(2.6)]
    rows_data=[
        ['1','New Verified','₩200,000','50 / month','50%','Business verified'],
        ['2','Trusted','₩750,000','200 / month','60%','90+ day clean record'],
        ['3','Premium / Dealer','₩2,000,000','1,000 / month','70%','High rating, no disputes'],
        ['4','Institutional','₩5,000,000+','Custom','70–75%','Audited partner status'],
    ]
    hdr_col=[BLU_H]*len(cols_h)
    xpos=ML
    for j,(ch,cw2) in enumerate(zip(cols_h,widths)):
        card(slide,xpos,t2,cw2,Inches(0.3),LIN_H,BLU_H,1.0)
        tb(slide,ch,xpos+Inches(0.05),t2+Inches(0.04),cw2,Inches(0.22),sz=8,bold=True,col=BLU_R)
        xpos+=cw2
    for ri,row in enumerate(rows_data):
        xpos=ML; col=_c([BLU_H,GRN_H,AMB_H,RED_H][ri])
        for ci,(cell,cw2) in enumerate(zip(row,widths)):
            bg=CARD_H if ri%2==0 else LIN_H
            card(slide,xpos,t2+Inches(0.32+ri*0.3),cw2,Inches(0.28),bg)
            tb(slide,cell,xpos+Inches(0.05),t2+Inches(0.35+ri*0.3),cw2,Inches(0.22),
               sz=8.5,col=col if ci==0 else TXT_R)
            xpos+=cw2

def s15_pool(slide,prem_img):
    slide_header(slide,"II.5.7 — PREMIUM POOL ALLOCATION")
    eyebrow(slide,'II.5.7 — PREMIUM POOL ALLOCATION')
    heading(slide,'Conservative Capital Management — Insurance ≠ Speculative Capital')
    subheading(slide,'Premiums are not free revenue. Every dollar is allocated with solvency as the first priority.')
    t=BODY_Y+Inches(0.1)
    img(slide,prem_img,Inches(0.1),t-Inches(0.1),Inches(6.0),Inches(5.7))
    x=ML+Inches(5.8); bw=CW-Inches(5.8)
    reserves=[
        ('Liquid Reserve  25%',     GRN_H,'Always liquid. Ineligible for DeFi yield. Pays claims immediately.'),
        ('Catastrophe Reserve  20%',BLU_H,'Used only during extreme events. Calibrated to Korean typhoon patterns.'),
        ('Safe-Rider Rebates  18%', CYN_H,'Returned quarterly to claim-free riders — the "fairness dividend."'),
        ('Expected Claims  12%',    RED_H,'Pays approved claims. Buffer above modeled ~9% loss ratio.'),
        ('Reinsurance  10%',        AMB_H,'External backstop. Required before launch. Covers tail risk above pool capacity.'),
        ('Treasury  10%',           MUT_H,'Protocol operational reserve and governance discretion.'),
        ('Protocol Fee  5%',        PUR_H,'Operating revenue captured for protocol development and team.'),
    ]
    for i,(label,col,desc) in enumerate(reserves):
        ry=t+Inches(i*0.75)
        tb(slide,'●',x,ry+Inches(0.07),Inches(0.22),Inches(0.28),sz=14,col=_c(col))
        tb(slide,label,x+Inches(0.25),ry+Inches(0.03),bw,Inches(0.28),sz=10,bold=True,col=_c(col))
        tb(slide,desc, x+Inches(0.25),ry+Inches(0.32),bw-Inches(0.1),Inches(0.35),sz=9,col=MUT_R)
    tb(slide,'The DeFi yield cap (up to 10% of excess reserves) is a ceiling — not a target. Only excess reserves above '
             'both the 25% liquidity floor and 20% catastrophe floor are eligible for deployment.',
       x,t+Inches(5.45),bw,Inches(0.52),sz=8.5,italic=True,col=AMB_R)

def s16_financials2(slide,pnl_img,lp_img):
    slide_header(slide,"II.5.8–5.9 — FX RISK & STRESS-TESTED FINANCIAL MODEL")
    eyebrow(slide,'II.5.8-5.9 — FX RISK & STRESS-TESTED FINANCIAL MODEL')
    heading(slide,'Commercially Viable Under Realistic Assumptions · Break-Even Year 3')
    t=BODY_Y
    # FX box top-left
    card(slide,ML,t,Inches(4.5),Inches(1.75),CARD_H,AMB_H,1.2)
    tb(slide,'II.5.8 — Currency & FX Risk',ML+Inches(0.18),t+Inches(0.12),Inches(4.1),Inches(0.3),
       sz=10,bold=True,col=AMB_R)
    tb(slide,'Premiums priced in ₩ (KRW) for riders but pool contributions held in USDC on-chain. '
             'Structural FX exposure managed via: (i) KRW-pegged stablecoins for a portion of reserves, '
             '(ii) USD-equivalent + KRW reserve ratios tracked in parallel, '
             '(iii) yield deployment paused when FX volatility exceeds governance thresholds.',
       ML+Inches(0.18),t+Inches(0.47),Inches(4.1),Inches(1.2),sz=9,col=TXT_R)
    # Stress params top-right
    card(slide,ML+Inches(4.7),t,CW-Inches(4.7),Inches(1.75),CARD_H,RED_H,1.2)
    tb(slide,'II.5.9 — Stress Parameters',ML+Inches(4.9),t+Inches(0.12),Inches(7.5),Inches(0.3),
       sz=10,bold=True,col=RED_R)
    stress=[
        ('Claim frequency','8% → 15%'),('Avg claim severity','₩1.2M → ₩2.5M'),
        ('Modeled loss ratio','~8.9% → ~34.7%'),('Safe-rider rebates','18% → ~10% (reduced under stress)'),
    ]
    for j,(label,val) in enumerate(stress):
        ry=t+Inches(0.52+j*0.28)
        sx=ML+Inches(4.9)
        tb(slide,f'{label}:',sx,ry,Inches(2.8),Inches(0.25),sz=9,col=MUT_R)
        tb(slide,val,sx+Inches(2.85),ry,Inches(4.5),Inches(0.25),sz=9,bold=True,col=RED_R)
    # P&L chart
    img(slide,pnl_img,ML,t+Inches(1.9),Inches(8.6),Inches(4.1))
    # LP chart
    img(slide,lp_img,ML+Inches(8.8),t+Inches(1.9),CW-Inches(8.8),Inches(4.1))
    tb(slide,'Total funding requirement: ~₩4.5B (~$3.3M USD) seed + Series A, separately from LP capital for pool underwriting.',
       ML,SH-Inches(0.45),CW,Inches(0.32),sz=9,italic=True,col=AMB_R,align=PP_ALIGN.CENTER)

def s17_competitive(slide,comp_img):
    slide_header(slide,"II.6.1 — COMPETITIVE ANALYSIS & MOAT")
    eyebrow(slide,'II.6.1 — COMPETITIVE ANALYSIS & MOAT')
    heading(slide,'Three Competitor Groups — One Intersection They Cannot Occupy')
    t=BODY_Y
    img(slide,comp_img,ML,t,CW,Inches(3.7))
    t2=t+Inches(3.85)
    moat=[
        ('Vertical Data\nNetwork Effect',BLU_H,
         'More bikes, riders, mechanics, dealers, partners → dataset value compounds for every participant. Structurally hard to replicate from a standing start.'),
        ('Dual-Protocol\nComposability',GRN_H,
         'MotoChain identity data enriches RideTrue behavior scoring. Each protocol makes the other more valuable. No single competitor can replicate both.'),
        ('Korea-First\nWedge',CYN_H,
         'Dense urban motorcycle market with realistic expansion to Vietnam, Indonesia, Thailand — Asia\'s largest two-wheeler markets.'),
        ('Consumer-Grade\nUX',AMB_H,
         'Account abstraction + gasless transactions target non-crypto users. Competitors in all three categories require crypto sophistication.'),
    ]
    mw=CW/4
    for i,(title,col,desc) in enumerate(moat):
        x=ML+i*mw
        card(slide,x,t2,mw-Inches(0.12),Inches(2.2),CARD_H,col,1.2)
        tb(slide,title,x+Inches(0.14),t2+Inches(0.12),mw-Inches(0.28),Inches(0.5),sz=10.5,bold=True,col=_c(col))
        tb(slide,desc, x+Inches(0.14),t2+Inches(0.68),mw-Inches(0.28),Inches(1.42),sz=9.5,col=TXT_R)

def s18_security(slide,sec_img):
    slide_header(slide,"II.6.2 — SECURITY ROADMAP")
    eyebrow(slide,'II.6.2 — SECURITY ROADMAP')
    heading(slide,'Security Is a Launch Gate — Not a Post-Launch Concern')
    subheading(slide,'The protocol will not accept user funds until both independent audits complete and all critical findings are resolved.')
    img(slide,sec_img,ML,BODY_Y+Inches(0.2),CW,Inches(3.5))
    t2=BODY_Y+Inches(3.85)
    opsec_items=[
        '5-of-9 multi-sig wallet with geographically distributed signers',
        '48-hour minimum time-lock on any non-emergency parameter change',
        'Emergency pause restricted to narrow scope with documented invocation criteria',
        'Public incident response runbook with quarterly drills',
        'Annual third-party penetration testing for mobile and telematics back-end systems',
        'Each material protocol upgrade preceded by an incremental audit',
    ]
    card(slide,ML,t2,CW,Inches(2.2),CARD_H,GRN_H,1.2)
    tb(slide,'Operational Security Commitments',ML+Inches(0.2),t2+Inches(0.12),CW,Inches(0.32),sz=11,bold=True,col=GRN_R)
    bxt=slide.shapes.add_textbox(ML+Inches(0.2),t2+Inches(0.5),CW-Inches(0.4),Inches(1.55))
    tf=bxt.text_frame; tf.word_wrap=True
    for j,item in enumerate(opsec_items):
        p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
        p.space_before=Pt(3)
        r=p.add_run(); r.text=f'• {item}'
        r.font.size=Pt(9.5); r.font.color.rgb=TXT_R

def s19_blackswan(slide,lp_img):
    slide_header(slide,"II.6.3–6.4 — ECONOMIC STRESS TESTING & LIQUIDITY SAFEGUARDS")
    eyebrow(slide,'II.6.3-6.4 — BLACK SWAN PLAN & LIQUIDITY SAFEGUARDS')
    heading(slide,'9 Safeguards Against Pool Collapse — Korea Typhoons Are Base-Case, Not Tail Risk')
    t=BODY_Y
    safeguards=[
        ('01','Mandatory\nLiquid Reserve','≥25% of pool always\nliquid — ineligible\nfor DeFi yield.',GRN_H),
        ('02','Catastrophe\nReserve Fund','20% of every premium\nrouted to dedicated\ncatastrophe fund.',BLU_H),
        ('03','Reinsurance\nBackstop','External partner covers\ntail risk above pool\ncapacity. Required pre-launch.',CYN_H),
        ('04','Emergency\nPayout Limits','Circuit breaker → 50%\nimmediate, 50% over\n3–6 months.',AMB_H),
        ('05','Pro-Rata\nDistribution','If undercollateralized,\nall claimants receive\nproportional share.',RED_H),
        ('06','Catastrophe\nCircuit Breaker','Triggers on: >20%\npool claims in 7 days,\nreserve ratio <15%.',BLU_H),
        ('07','Conservative\nDeFi Yield','Only excess reserves\ndeployed. Cap 10–15%.\nBlue-chip only.',GRN_H),
        ('08','Dynamic\nPremium Adj.','Post-shock: +10–20%\ntemporary increase to\nrebuild reserves.',AMB_H),
        ('09','Coverage\nExclusions & Caps','Max normal payout:\n₩5M. Max catastrophe:\n₩3M immediate.',CYN_H),
    ]
    # 3×3 grid
    gw=(CW-Inches(0.3))/3; gh=Inches(1.45)
    for i,(num,title,desc,col) in enumerate(safeguards):
        r=i//3; c=i%3
        x=ML+c*(gw+Inches(0.15)); y=t+r*(gh+Inches(0.12))
        card(slide,x,y,gw,gh,CARD_H,col,1.2)
        tb(slide,num,  x+Inches(0.12),y+Inches(0.08),Inches(0.3),Inches(0.3),sz=9,bold=True,col=_c(col))
        tb(slide,title,x+Inches(0.42),y+Inches(0.08),gw-Inches(0.55),Inches(0.5),sz=10,bold=True,col=_c(col))
        tb(slide,desc, x+Inches(0.12),y+Inches(0.62),gw-Inches(0.25),Inches(0.75),sz=8.5,col=TXT_R)
    tb(slide,'Liquidity safeguards: 25% liquid floor · capped DeFi deployment · 30-day LP withdrawal notice · '
             'pro-rata exit queue in emergency mode · stablecoin diversification (USDC + KRW-pegged) · real-time public liquidity dashboard.',
       ML,SH-Inches(0.48),CW,Inches(0.38),sz=9,italic=True,col=AMB_R,align=PP_ALIGN.CENTER)

def s20_oracle_fraud(slide,oracle_img):
    slide_header(slide,"II.6.5–6.6 — ORACLE RISK & MECHANIC/VALIDATOR FRAUD")
    eyebrow(slide,'II.6.5-6.6 — ORACLE RISK & MECHANIC FRAUD')
    heading(slide,'The Last Line of Defense: Data Integrity End to End')
    t=BODY_Y
    img(slide,oracle_img,ML,t,Inches(7.5),Inches(4.0))
    x2=ML+Inches(7.7); bw2=CW-Inches(7.7)
    tb(slide,'Oracle Risk Mitigation',x2,t,bw2,Inches(0.32),sz=12,bold=True,col=CYN_R)
    oracle_pts=[
        'Multi-source aggregation — theft claim needs police record + vehicle registry + insurer confirmation',
        'No single oracle trusted in isolation — consensus threshold required',
        'Stale / conflicting data → verification-pending state, not auto-execution',
        'High-value claims require higher consensus thresholds + longer review windows',
    ]
    bxt=slide.shapes.add_textbox(x2,t+Inches(0.38),bw2,Inches(2.0))
    tf=bxt.text_frame; tf.word_wrap=True
    for j,pt in enumerate(oracle_pts):
        p=tf.paragraphs[0] if j==0 else tf.add_paragraph()
        p.space_before=Pt(7); r=p.add_run(); r.text=f'• {pt}'
        r.font.size=Pt(9.5); r.font.color.rgb=TXT_R
    tb(slide,'Mechanic Fraud Summary',x2,t+Inches(2.6),bw2,Inches(0.32),sz=12,bold=True,col=AMB_R)
    fraud_layers=[
        ('Prevention',GRN_H,'Co-signing + hardware telemetry eliminate categories of unilateral fraud at the source.'),
        ('Detection', AMB_H,'Bounty market + monetization-triggered re-verification surface fraud beyond the escrow window.'),
        ('Enforcement',RED_H,'Tiered staking + risk-tiered escrow + per-entry stake-at-risk make detected fraud economically irrational.'),
    ]
    for j,(layer,col,desc) in enumerate(fraud_layers):
        ry=t+Inches(3.0+j*0.85)
        card(slide,x2,ry,bw2,Inches(0.72),CARD_H,col,1.0)
        tb(slide,layer,x2+Inches(0.12),ry+Inches(0.08),bw2,Inches(0.28),sz=10,bold=True,col=_c(col))
        tb(slide,desc, x2+Inches(0.12),ry+Inches(0.38),bw2-Inches(0.2),Inches(0.28),sz=8.5,col=TXT_R)
    # Summary bar
    card(slide,ML,t+Inches(4.2),Inches(7.5),Inches(0.95),CARD_H,GRN_H,1.2)
    tb(slide,'Design Goal:  The architecture does not assume mechanics are honest or dishonest by default. '
             'It assumes they are rational actors — and keeps the expected value of fraud negative across all attack vectors.',
       ML+Inches(0.2),t+Inches(4.32),Inches(7.1),Inches(0.68),sz=10,italic=True,col=TXT_R)

# ── ASSEMBLE ───────────────────────────────────────────────────────────────
print('Generating charts …')
vin_img    = ch_vin()
zk_img     = ch_zk()
claim_img  = ch_claim()
market_img = ch_market_bars()
rev_img    = ch_revenue_pie()
tok_img    = ch_token_donut()
prem_img   = ch_premium_pie()
mech_img   = ch_mechanic_defense()
pnl_img    = ch_pnl()
lp_img     = ch_lp()
comp_img   = ch_competitive()
sec_img    = ch_security()
oracle_img = ch_oracle()
ts_img     = ch_two_sided()
print('  ✓ all 14 charts generated')

TOTAL=20
prs=new_prs()

plan=[
    ('Title',                        lambda s: s01_title(s)),
    ('I. Elevator Pitch',            lambda s: s02_pitch(s)),
    ('Mission Statement',            lambda s: s03_mission(s)),
    ('II.1 Problem — Market Scale',  lambda s: s04_problem_market(s,market_img)),
    ('II.1 Why Blockchain',          lambda s: s05_why_blockchain(s)),
    ('II.2 VIN Lifecycle',           lambda s: s06_architecture(s,vin_img)),
    ('II.2.1 ZK + Anti-Spoofing',    lambda s: s07_zk_antispoofing(s,zk_img)),
    ('II.2.2 Claims & Oracles',      lambda s: s08_claims_oracles(s,claim_img,oracle_img)),
    ('II.3 User Experience',         lambda s: s09_ux(s)),
    ('II.4 Target Audience',         lambda s: s10_audience(s,ts_img)),
    ('II.5.1-2 Revenue Streams',     lambda s: s11_revenue(s,rev_img)),
    ('II.5.3-4 Financial Model',     lambda s: s12_financials(s)),
    ('II.5.5 Token Design',          lambda s: s13_token(s,tok_img)),
    ('II.5.6 Mechanic Trust',        lambda s: s14_mechanic(s,mech_img)),
    ('II.5.7 Premium Pool',          lambda s: s15_pool(s,prem_img)),
    ('II.5.8-9 FX + P&L',           lambda s: s16_financials2(s,pnl_img,lp_img)),
    ('II.6.1 Competitive + Moat',    lambda s: s17_competitive(s,comp_img)),
    ('II.6.2 Security Roadmap',      lambda s: s18_security(s,sec_img)),
    ('II.6.3-4 Black Swan + Liquidity',lambda s: s19_blackswan(s,lp_img)),
    ('II.6.5-6 Oracle + Fraud',      lambda s: s20_oracle_fraud(s,oracle_img)),
]

print('Building slides …')
for n,(name,builder) in enumerate(plan,1):
    print(f'  {n:2d}/{TOTAL}  {name}')
    sl=blank(prs)
    builder(sl)
    if n>1:
        pagenum(sl,n,TOTAL)

out='0522_G1_Interim_Report.pptx'
prs.save(out)
print(f'\n✓  Saved → {out}  ({TOTAL} slides)')
